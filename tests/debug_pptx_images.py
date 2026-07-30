"""
debug_pptx_images.py
─────────────────────────────────────────────
Inspecte toutes les images d'un PPTX pour trouver
celle qui cause le nan/inf dans Qwen.

Usage :
  python debug_pptx_images.py "PRESENTATION SOUTENANCE.pptx"
"""
import sys
import io
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_pptx_images(path: str):
    prs  = Presentation(path)
    total, ok, bad = 0, 0, 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "image"):
                continue
            total += 1
            try:
                blob         = shape.image.blob
                content_type = shape.image.content_type
                size_kb      = len(blob) / 1024

                # Test 1 : PIL peut-il décoder ?
                img = Image.open(io.BytesIO(blob))
                img.load()
                w, h = img.size
                mode = img.mode

                # Test 2 : dimensions suspectes ?
                warn = ""
                if w < 10 or h < 10:
                    warn = "⚠️  TROP PETITE"
                elif w > 4000 or h > 4000:
                    warn = "⚠️  TRÈS GRANDE"
                elif w * h == 0:
                    warn = "❌ DIMENSION ZÉRO"

                print(f"  Slide {slide_num} shape {shape_idx} : "
                      f"{w}×{h} {mode} | {content_type} | {size_kb:.1f}KB {warn}")
                ok += 1

            except Exception as exc:
                bad += 1
                try:
                    ct = shape.image.content_type
                except Exception:
                    ct = "inconnu"
                print(f"  ❌ Slide {slide_num} shape {shape_idx} : "
                      f"ERREUR DÉCODAGE | type={ct} | {exc}")

    print(f"\nTotal : {total} images | OK : {ok} | Problèmes : {bad}")
    if bad == 0 and total > 0:
        print("\n→ Toutes les images se décodent. Le problème vient")
        print("  probablement d'une image avec des dimensions inhabituelles")
        print("  (ex: très grande résolution). Vérifiez les '⚠️' ci-dessus.")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage : python debug_pptx_images.py <fichier.pptx>")
        sys.exit(1)
    inspect_pptx_images(sys.argv[1])