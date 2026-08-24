# -*- coding: utf-8 -*-
from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Mapping, Optional

from .cleaner import clean_text
from .origin_detector import infer_origin
from .document_type_classifier import enrich_document_type


def _load_supported_extensions_from_extraction_router() -> set[str]:
    """
    Liste unique des extensions acceptées par le loader NLP.

    Objectif :
    - utiliser la même couverture que modules/extraction/router.py ;
    - éviter que le backend accepte un fichier mais que document_loader le rejette ;
    - garder un fallback complet si le router n'est pas importable.
    """
    fallback = {
        # PDF
        ".pdf",

        # Office / documents
        ".docx", ".doc", ".docm",
        ".pptx", ".ppt", ".pptm",
        ".xlsx", ".xls", ".xlsm", ".csv",

        # Texte / données
        ".txt", ".md", ".json",

        # Emails
        ".eml", ".msg",

        # Images
        ".png", ".jpg", ".jpeg", ".tiff", ".tif",
        ".bmp", ".gif", ".webp", ".svg",

        # Audio
        ".mp3", ".wav", ".m4a", ".aac", ".flac",
        ".ogg", ".opus", ".wma",

        # Vidéo
        ".mp4", ".mov", ".avi", ".mkv", ".webm",
        ".mpeg", ".mpg", ".3gp",
    }

    try:
        from modules.extraction.router import EXTENSION_MAP, AUDIO_VIDEO_EXTENSIONS

        exts = set(EXTENSION_MAP.keys()) | set(AUDIO_VIDEO_EXTENSIONS)

        # Extensions utiles côté NLP/debug.
        exts |= {".txt", ".md", ".json"}

        # Formats macro Office. Le router peut les détecter par magic bytes,
        # mais on les accepte explicitement ici.
        exts |= {".docm", ".pptm"}

        return {str(ext).lower() for ext in exts if str(ext).startswith(".")}

    except Exception:
        return fallback


SUPPORTED_EXTENSIONS = _load_supported_extensions_from_extraction_router()


def normalize_chunks(value: Any) -> List[str]:
    if value is None:
        return []

    if isinstance(value, str):
        return [value]

    out: List[str] = []

    if isinstance(value, list):
        for item in value:
            if isinstance(item, str):
                out.append(item)

            elif isinstance(item, dict):
                for key in ["text", "content", "chunk", "page_text", "rag_chunk"]:
                    if isinstance(item.get(key), str):
                        out.append(item[key])
                        break

            elif hasattr(item, "text") and isinstance(getattr(item, "text"), str):
                out.append(getattr(item, "text"))

            elif hasattr(item, "content") and isinstance(getattr(item, "content"), str):
                out.append(getattr(item, "content"))

    elif isinstance(value, dict):
        for key in ["text", "content", "chunk", "page_text", "rag_chunk"]:
            if isinstance(value.get(key), str):
                out.append(value[key])

    return out


def extract_text_from_result(result: Any) -> str:
    """
    Convertit un résultat ExtractionResult / dict / objet en texte brut stable.
    """
    if result is None:
        return ""

    chunks: List[str] = []

    if isinstance(result, dict):
        for key in ["text_chunks", "chunks", "pages", "texts", "content"]:
            chunks.extend(normalize_chunks(result.get(key)))

        if isinstance(result.get("text"), str):
            chunks.append(result["text"])

        # Certains extracteurs mettent la description image ici.
        chunks.extend(normalize_chunks(result.get("visual_chunks")))

    else:
        for attr in ["text_chunks", "chunks", "pages", "texts", "content"]:
            if hasattr(result, attr):
                chunks.extend(normalize_chunks(getattr(result, attr)))

        if hasattr(result, "text") and isinstance(getattr(result, "text"), str):
            chunks.append(getattr(result, "text"))

        if hasattr(result, "visual_chunks"):
            chunks.extend(normalize_chunks(getattr(result, "visual_chunks")))

    if isinstance(result, str):
        chunks.append(result)

    return "\n\n".join(
        chunk.strip()
        for chunk in chunks
        if isinstance(chunk, str) and chunk.strip()
    )


def extract_with_ennosmart_router(
    path: str,
    *,
    vision_mode: Optional[str] = None,
    formula_mode: Optional[str] = None,
) -> Optional[str]:
    """
    Extraction prioritaire via modules.extraction.router.

    Paramètres pilotables par variables d'environnement :
    - ENNOSMART_EXTRACTION_VISION_MODE = text_only | auto | fast | full
    - ENNOSMART_EXTRACTION_FORMULA_MODE = off | fast | explain
    - ENNOSMART_ENABLE_TRANSCRIPTION = 1 | 0
    """
    try:
        from modules.extraction.router import extract
    except Exception:
        return None

    try:
        result = extract(
            path,
            vision_mode=(
                vision_mode
                or os.getenv("ENNOSMART_EXTRACTION_VISION_MODE", "text_only")
            ),
            formula_mode=(
                formula_mode
                or os.getenv("ENNOSMART_EXTRACTION_FORMULA_MODE", "off")
            ),
            enable_transcription=os.getenv("ENNOSMART_ENABLE_TRANSCRIPTION", "1").strip() != "0",
        )
    except Exception:
        return None

    text = extract_text_from_result(result)
    text = clean_text(text)

    return text if text.strip() else None


def fallback_extract(path: str) -> str:
    """
    Fallback léger si le router EnnoSmart échoue.

    Ce fallback ne remplace pas le router :
    - il sert seulement à éviter qu'un document soit totalement perdu ;
    - pour audio/vidéo/image pure, le texte peut rester vide si vision/transcription indisponible.
    """
    ext = Path(path).suffix.lower()

    if ext in {".txt", ".md", ".csv"}:
        return read_text_file(path)

    if ext == ".json":
        return read_json_file(path)

    if ext in {".docx", ".docm"}:
        return read_docx(path)

    if ext in {".pptx", ".pptm"}:
        return read_pptx(path)

    if ext in {".xlsx", ".xls", ".xlsm"}:
        return read_excel(path)

    if ext == ".pdf":
        return read_pdf(path)

    if ext == ".eml":
        return read_eml(path)

    if ext == ".msg":
        return read_msg(path)

    # Anciens formats Office binaires : le router ou LibreOffice doit les gérer.
    # Ici on évite une erreur bloquante.
    return ""


def read_text_file(path: str) -> str:
    for enc in ["utf-8", "utf-8-sig", "cp1252", "latin-1"]:
        try:
            return Path(path).read_text(encoding=enc, errors="ignore")
        except Exception:
            pass

    return ""


def read_json_file(path: str) -> str:
    raw = read_text_file(path)

    try:
        data = json.loads(raw)
    except Exception:
        return raw

    parts: List[str] = []

    def walk(value: Any):
        if isinstance(value, dict):
            for key, child in value.items():
                if str(key).lower() in {
                    "source_path",
                    "file_category",
                    "visual_chunks",
                    "structured_data",
                    "attachments_paths",
                    "extraction_errors",
                }:
                    continue
                walk(child)

        elif isinstance(value, list):
            for item in value:
                walk(item)

        elif isinstance(value, str) and len(value.strip()) > 25:
            parts.append(value.strip())

    walk(data)

    return "\n\n".join(parts)


def read_docx(path: str) -> str:
    try:
        from docx import Document

        doc = Document(path)
        parts: List[str] = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            for row in table.rows:
                cells = [
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text and cell.text.strip()
                ]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n\n".join(parts)

    except Exception:
        return ""


def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        parts: List[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts: List[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_parts.append(text)

            if slide_parts:
                parts.append(f"[SLIDE {i}]\n" + "\n".join(slide_parts))

        return "\n\n".join(parts)

    except Exception:
        return ""


def read_excel(path: str) -> str:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        parts: List[str] = []

        for ws in wb.worksheets:
            rows: List[str] = []

            for row in ws.iter_rows(values_only=True):
                vals = [
                    str(value).strip()
                    for value in row
                    if value is not None and str(value).strip()
                ]

                if vals:
                    rows.append(" | ".join(vals))

            if rows:
                parts.append(f"[FEUILLE : {ws.title}]\n" + "\n".join(rows[:350]))

        return "\n\n".join(parts)

    except Exception:
        return ""


def read_pdf(path: str) -> str:
    try:
        try:
            from pypdf import PdfReader
        except Exception:
            from PyPDF2 import PdfReader

        reader = PdfReader(path)
        parts: List[str] = []

        for i, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()

            if text:
                parts.append(f"[PAGE {i}]\n{text}")

        return "\n\n".join(parts)

    except Exception:
        return ""


def read_eml(path: str) -> str:
    try:
        from email import policy
        from email.parser import BytesParser

        with open(path, "rb") as f:
            msg = BytesParser(policy=policy.default).parse(f)

        parts: List[str] = []

        for label, key in [("Objet", "subject"), ("De", "from"), ("Date", "date")]:
            if msg.get(key):
                parts.append(f"{label} : {msg.get(key)}")

        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    try:
                        content = part.get_content()
                    except Exception:
                        content = ""
                    if content:
                        parts.append(content)

        elif msg.get_content_type() == "text/plain":
            parts.append(msg.get_content())

        return "\n\n".join(parts)

    except Exception:
        return ""


def read_msg(path: str) -> str:
    try:
        import extract_msg

        msg = extract_msg.Message(path)
        parts: List[str] = []

        try:
            if msg.subject:
                parts.append(f"Objet : {msg.subject}")

            if msg.sender:
                parts.append(f"De : {msg.sender}")

            if msg.date:
                parts.append(f"Date : {msg.date}")

            if msg.body:
                parts.append(msg.body)

        finally:
            try:
                msg.close()
            except Exception:
                pass

        return "\n\n".join(parts)

    except Exception:
        return ""


def load_document(
    path: str,
    use_ennosmart_extraction: bool = True,
    *,
    metadata: Optional[Mapping[str, Any]] = None,
    vision_mode: Optional[str] = None,
    formula_mode: Optional[str] = None,
) -> Dict[str, Any]:
    p = Path(path)
    ext = p.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        return {
            "document": p.name,
            "source_path": str(p),
            "extension": ext,
            "text": "",
            "loader": "unsupported",
            "chars": 0,
            "error": f"Extension non supportée : {ext}",
        }

    text = (
        extract_with_ennosmart_router(
            str(p),
            vision_mode=vision_mode,
            formula_mode=formula_mode,
        )
        if use_ennosmart_extraction
        else None
    )
    loader = "modules.extraction.router" if text else "fallback"

    if not text:
        text = fallback_extract(str(p))

    text = clean_text(text)

    origin = infer_origin(p.name, text)

    doc = {
        "document": p.name,
        "source_path": str(p),
        "extension": ext,
        "text": text,
        "loader": loader,
        "chars": len(text or ""),
        "error": None if text else "Aucun texte extrait",
        **origin,
    }

    doc = enrich_document_type(doc)

    # Les métadonnées déclarées par l'application sont la source de vérité sur
    # l'usage du fichier. La classification automatique reste traçable, mais
    # elle ne doit pas exclure un document explicitement déposé comme preuve
    # courante, même si son contenu ressemble à un CIR déjà rédigé.
    declared = dict(metadata or {})
    if declared:
        auto_document_type = doc.get("document_type")
        auto_content_origin = doc.get("content_origin")
        doc.update(declared)
        doc["auto_document_type"] = auto_document_type
        doc["auto_content_origin"] = auto_content_origin

        if str(declared.get("declared_mode") or "").strip().lower() == "raw":
            doc.update(
                {
                    "content_origin": "raw_client_document",
                    "source_policy": "core_or_useful",
                    "current_project_evidence": True,
                    "declared_raw_document": True,
                    "cir_final_validated": False,
                    "not_final_cir": True,
                }
            )
            doc["document_weight"] = max(
                float(doc.get("document_weight") or 0.0),
                1.0,
            )
            doc["source_weight"] = max(
                float(doc.get("source_weight") or 0.0),
                float(doc["document_weight"]),
            )

    return doc


def load_documents(
    paths: List[str],
    use_ennosmart_extraction: bool = True,
    include_cir_final: bool = False,
    *,
    metadata_by_path: Optional[Mapping[str, Mapping[str, Any]]] = None,
    vision_mode: Optional[str] = None,
    formula_mode: Optional[str] = None,
) -> List[Dict[str, Any]]:
    docs: List[Dict[str, Any]] = []
    metadata_lookup = {
        str(Path(key).resolve()).lower(): dict(value)
        for key, value in (metadata_by_path or {}).items()
    }

    for path in paths:
        d = load_document(
            path,
            use_ennosmart_extraction=use_ennosmart_extraction,
            metadata=metadata_lookup.get(str(Path(path).resolve()).lower()),
            vision_mode=vision_mode,
            formula_mode=formula_mode,
        )

        if not d.get("text", "").strip():
            # On garde le comportement actuel : les documents sans texte exploitable
            # ne partent pas au NLP.
            continue

        if d.get("content_origin") == "cir_final" and not include_cir_final:
            d["skipped_reason"] = "cir_final_excluded"
            continue

        docs.append(d)

    return docs
