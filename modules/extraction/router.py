"""
modules/extraction/router.py — EnnoSmart extraction router optimisé

Nouveautés :
- formula_mode = "off" | "fast" | "explain"
  Par défaut : "off" pendant la refonte NLP evidence-first.

- vision_mode = "text_only" | "auto" | "fast" | "full"
  Par défaut : "text_only" pour éviter le coût vision pendant les tests NLP.

- Transcription audio/vidéo via faster-whisper :
  MP3/WAV/M4A/MP4/MOV/etc. → text_chunks compatibles NLP + RAG.

- Regroupement intelligent des segments audio :
  évite d'indexer 900+ petits chunks pour un seul audio long.

- Warmup vision seulement après filtrage réel des images utiles.
- Limite images par document pour éviter des traitements trop longs.
- Compatible avec l'ancien paramètre enable_formulas=True/False.
"""

from __future__ import annotations

import logging
import mimetypes
import re
import tempfile
import zipfile
from pathlib import Path
from typing import Optional

from modules.extraction.base import ExtractionResult, FileCategory, SourceTag

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────────────────────────
# Extensions supportées
# ──────────────────────────────────────────────────────────────────────────────

EXTENSION_MAP: dict[str, FileCategory] = {
    ".pdf": FileCategory.PDF_NATIVE,

    ".docx": FileCategory.DOCX,
    ".docm": FileCategory.DOCX,
    ".doc": FileCategory.DOCX,

    ".pptx": FileCategory.PPTX,
    ".pptm": FileCategory.PPTX,
    ".ppt": FileCategory.PPTX,

    ".eml": FileCategory.EMAIL,
    ".msg": FileCategory.EMAIL,

    ".xlsx": FileCategory.EXCEL,
    ".xlsm": FileCategory.EXCEL,
    ".xls": FileCategory.EXCEL,
    ".csv": FileCategory.EXCEL,

    ".png": FileCategory.IMAGE,
    ".jpg": FileCategory.IMAGE,
    ".jpeg": FileCategory.IMAGE,
    ".tiff": FileCategory.IMAGE,
    ".tif": FileCategory.IMAGE,
    ".bmp": FileCategory.IMAGE,
    ".gif": FileCategory.IMAGE,
    ".webp": FileCategory.IMAGE,
    ".svg": FileCategory.IMAGE,
}

AUDIO_EXTENSIONS: set[str] = {
    ".mp3",
    ".wav",
    ".m4a",
    ".aac",
    ".flac",
    ".ogg",
    ".opus",
    ".wma",
}

VIDEO_EXTENSIONS: set[str] = {
    ".mp4",
    ".mov",
    ".avi",
    ".mkv",
    ".webm",
    ".mpeg",
    ".mpg",
    ".3gp",
}

AUDIO_VIDEO_EXTENSIONS: set[str] = AUDIO_EXTENSIONS | VIDEO_EXTENSIONS

# Extensions Office Open XML : on lit le header ZIP pour distinguer Word/PPT/Excel.
OPENXML_EXTENSIONS: set[str] = {
    ".docx",
    ".docm",
    ".pptx",
    ".pptm",
    ".xlsx",
    ".xlsm",
}

# Archives : le ZIP est supporté nativement. RAR/7Z peuvent être ajoutés plus tard.
ARCHIVE_EXTENSIONS: set[str] = {".zip"}

MAX_ARCHIVE_FILES = 80
MAX_ARCHIVE_TOTAL_BYTES = 250 * 1024 * 1024  # 250 MB
MAX_ARCHIVE_SINGLE_FILE_BYTES = 100 * 1024 * 1024  # 100 MB

MAGIC_BYTES = [
    (b"%PDF", FileCategory.PDF_NATIVE),
    (b"PK\x03\x04", FileCategory.DOCX),
    (b"\xd0\xcf\x11\xe0", FileCategory.DOCX),
]

MAX_IMAGES_PER_SLIDE = 3
MAX_IMAGES_PER_DOCUMENT_AUTO = 2
MAX_IMAGES_PER_DOCUMENT_FAST = 3
MAX_IMAGES_PER_DOCUMENT_FULL = 3

DEFAULT_TRANSCRIPTION_CHUNK_SECONDS = 90
DEFAULT_TRANSCRIPTION_CHUNK_MAX_CHARS = 2500


# ──────────────────────────────────────────────────────────────────────────────
# Détection type fichier
# ──────────────────────────────────────────────────────────────────────────────

def _is_audio_video(path: Path) -> bool:
    return path.suffix.lower() in AUDIO_VIDEO_EXTENSIONS


def _is_video(path: Path) -> bool:
    return path.suffix.lower() in VIDEO_EXTENSIONS


def _audio_file_category() -> FileCategory:
    """
    Garde le router compatible même si ton enum FileCategory
    n'a pas encore AUDIO_VIDEO.
    """
    for name in ("AUDIO_VIDEO", "AUDIO", "VIDEO", "TRANSCRIPTION"):
        if hasattr(FileCategory, name):
            return getattr(FileCategory, name)

    return FileCategory.UNKNOWN



def _archive_file_category() -> FileCategory:
    """
    Garde le router compatible même si FileCategory n'a pas ARCHIVE.
    """
    for name in ("ARCHIVE", "ZIP", "CONTAINER"):
        if hasattr(FileCategory, name):
            return getattr(FileCategory, name)

    return FileCategory.UNKNOWN


def _detect_category(path: Path) -> FileCategory:
    ext = path.suffix.lower()

    if ext in ARCHIVE_EXTENSIONS:
        return _archive_file_category()

    if ext in AUDIO_VIDEO_EXTENSIONS:
        return _audio_file_category()

    if ext in EXTENSION_MAP and ext not in OPENXML_EXTENSIONS:
        return EXTENSION_MAP[ext]

    try:
        with open(path, "rb") as f:
            header = f.read(8)

        for magic, cat in MAGIC_BYTES:
            if header.startswith(magic):
                if magic == b"PK\x03\x04":
                    return _detect_zip_subtype(path, ext)

                if magic == b"\xd0\xcf\x11\xe0":
                    return FileCategory.EMAIL if ext == ".msg" else FileCategory.DOCX

                return cat

    except Exception:
        pass

    if ext in EXTENSION_MAP:
        return EXTENSION_MAP[ext]

    mime, _ = mimetypes.guess_type(str(path))

    if mime:
        for k, v in [
            ("pdf", FileCategory.PDF_NATIVE),
            ("word", FileCategory.DOCX),
            ("presentation", FileCategory.PPTX),
            ("spreadsheet", FileCategory.EXCEL),
            ("image", FileCategory.IMAGE),
            ("message", FileCategory.EMAIL),
        ]:
            if k in mime:
                return v

        if mime.startswith("audio/") or mime.startswith("video/"):
            return _audio_file_category()

    return FileCategory.UNKNOWN


def _detect_zip_subtype(path: Path, ext: str) -> FileCategory:
    try:
        import zipfile

        with zipfile.ZipFile(str(path)) as zf:
            names = zf.namelist()

            if any("word/" in n for n in names):
                return FileCategory.DOCX

            if any("ppt/" in n for n in names):
                return FileCategory.PPTX

            if any("xl/" in n for n in names):
                return FileCategory.EXCEL

    except Exception:
        pass

    return EXTENSION_MAP.get(ext, FileCategory.UNKNOWN)


def _category_label(category: FileCategory, path: Path) -> str:
    if _is_audio_video(path):
        return "audio_video"

    return getattr(category, "value", str(category))


def _normalize_formula_mode(
    formula_mode: str | None = None,
    enable_formulas: Optional[bool] = None,
) -> str:
    if enable_formulas is False:
        return "off"

    mode = (formula_mode or "off").lower().strip()

    if mode in {"false", "none", "disabled", "no", "0"}:
        return "off"

    if mode not in {"off", "fast", "explain"}:
        logger.warning("formula_mode inconnu %r → off", formula_mode)
        return "off"

    return mode


def _normalize_vision_mode(mode: str | None) -> str:
    m = (mode or "text_only").lower().strip()

    if m in {"none", "off", "no_vision", "disabled"}:
        return "text_only"

    if m not in {"text_only", "auto", "fast", "full"}:
        logger.warning("vision_mode inconnu %r → text_only", mode)
        return "text_only"

    return m


# ──────────────────────────────────────────────────────────────────────────────
# Utilitaires chunks
# ──────────────────────────────────────────────────────────────────────────────

def _slide_num(header: str) -> Optional[int]:
    m = re.search(r"(?:SLIDE|PAGE)\s+(\d+)", str(header or ""), re.IGNORECASE)
    return int(m.group(1)) if m else None


def _build_chunk_index(chunks: list[str]) -> dict[int, int]:
    idx: dict[int, int] = {}

    for i, c in enumerate(chunks or []):
        n = _slide_num(c.split("\n")[0] if c else "")

        if n is not None:
            idx[n] = i

    return idx


def _safe_normalize_image(img_bytes: bytes) -> Optional[bytes]:
    try:
        from PIL import Image
        import io

        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        img.load()

        w, h = img.size

        if w < 50 or h < 50:
            return None

        aspect = w / h if h > 0 else 1

        if aspect > 5 or aspect < 0.2:
            return None

        if w == 1135 and h == 289:
            return None

        canvas = Image.new("RGB", (448, 448), (255, 255, 255))
        img.thumbnail((448, 448), Image.LANCZOS)
        canvas.paste(img, ((448 - img.width) // 2, (448 - img.height) // 2))

        buf = io.BytesIO()
        canvas.save(buf, "PNG")

        return buf.getvalue()

    except Exception:
        return None


# ──────────────────────────────────────────────────────────────────────────────
# Formules
# ──────────────────────────────────────────────────────────────────────────────

def _inject_formulas(
    chunks: list[str],
    context_file: str = "",
    formula_mode: str = "off",
) -> list[str]:
    mode = _normalize_formula_mode(formula_mode)

    if mode == "off":
        return chunks

    try:
        from modules.extraction.formula.formula import extract_formulas, _detect
    except Exception as exc:
        logger.debug("Module formule indisponible : %s", exc)
        return chunks

    result: list[str] = []

    for chunk in chunks or []:
        num = _slide_num(chunk.split("\n")[0] if chunk else "")
        ctx = (f"Slide/Page {num} — " if num else "") + context_file

        has, _ = _detect(chunk, ctx)

        if not has:
            result.append(chunk)
            continue

        fmls = extract_formulas(
            text=chunk,
            context=ctx,
            formula_mode=mode,
        )

        if not fmls:
            result.append(chunk)
            continue

        seen: set[str] = set()
        section = "\n[FORMULES DÉTECTÉES]\n"

        for i, f in enumerate(fmls, 1):
            k = re.sub(r"\s+", "", f.latex.strip().lower())

            if not k or k in seen:
                continue

            seen.add(k)

            section += (
                f"  {i}. LaTeX: {f.latex}\n"
                f"     Source: {f.source.value} | "
                f"Domaine: {f.domain.value} | "
                f"Confiance: {f.confidence:.2f}\n"
            )

        result.append(chunk.rstrip() + section if seen else chunk)

    return result


def _run_docx_omml(
    path: Path,
    chunks: list[str],
    formula_mode: str = "off",
) -> list[str]:
    if _normalize_formula_mode(formula_mode) == "off":
        return chunks

    try:
        import zipfile
        import xml.etree.ElementTree as ET
        from modules.extraction.formula.formula import extract_formulas

        ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"

        with zipfile.ZipFile(str(path)) as zf:
            if "word/document.xml" not in zf.namelist():
                return chunks

            root = ET.fromstring(zf.read("word/document.xml"))
            omaths = list(root.iter("{" + ns + "}oMath"))

        if not omaths:
            return chunks

        enriched = list(chunks)
        seen: set[str] = set()
        section = "\n[FORMULES OMML]\n"

        for omath in omaths:
            for r in extract_formulas(
                omml_xml=ET.tostring(omath, encoding="unicode"),
                context=path.name,
                formula_mode=formula_mode,
            ):
                k = re.sub(r"\s+", "", r.latex.strip().lower())

                if k and k not in seen:
                    seen.add(k)
                    section += f"  • {r.latex} — {r.source.value} | {r.confidence:.2f}\n"

        if enriched and seen:
            enriched[-1] = enriched[-1].rstrip() + section

        return enriched

    except Exception as exc:
        logger.debug("OMML DOCX : %s", exc)
        return chunks


def _run_pptx_omml(
    path: Path,
    chunks: list[str],
    formula_mode: str = "off",
) -> list[str]:
    if _normalize_formula_mode(formula_mode) == "off":
        return chunks

    try:
        import zipfile
        import xml.etree.ElementTree as ET
        from modules.extraction.formula.formula import extract_formulas

        ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        cidx = _build_chunk_index(chunks)
        enriched = list(chunks)

        with zipfile.ZipFile(str(path)) as zf:
            slide_files = sorted(
                f
                for f in zf.namelist()
                if f.startswith("ppt/slides/slide") and f.endswith(".xml")
            )

            for sf in slide_files:
                try:
                    tree = ET.fromstring(zf.read(sf))
                    omaths = list(tree.iter("{" + ns + "}oMath"))

                    if not omaths:
                        continue

                    m = re.search(r"slide(\d+)\.xml", sf)
                    snum = int(m.group(1)) if m else 0

                    ctx = (
                        enriched[cidx[snum]][:400]
                        if snum in cidx
                        else f"Slide {snum} R&D"
                    )

                    seen: set[str] = set()
                    section = ""

                    for omath in omaths:
                        for r in extract_formulas(
                            omml_xml=ET.tostring(omath, encoding="unicode"),
                            context=ctx,
                            formula_mode=formula_mode,
                        ):
                            k = re.sub(r"\s+", "", r.latex.strip().lower())

                            if k and k not in seen:
                                seen.add(k)
                                section += f"  • {r.latex} — {r.source.value} | {r.confidence:.2f}\n"

                    if section and snum in cidx:
                        enriched[cidx[snum]] = (
                            enriched[cidx[snum]].rstrip()
                            + "\n[FORMULES OMML]\n"
                            + section
                        )

                except Exception as exc:
                    logger.debug("OMML slide %s : %s", sf, exc)

        return enriched

    except Exception as exc:
        logger.debug("OMML PPTX : %s", exc)
        return chunks


# ──────────────────────────────────────────────────────────────────────────────
# Images / Vision
# ──────────────────────────────────────────────────────────────────────────────

def _limit_images_per_slide(items: list[dict]) -> list[dict]:
    from collections import defaultdict
    import io

    groups: dict = defaultdict(list)

    for item in items or []:
        groups[item.get("slide") or item.get("page") or 0].append(item)

    result = []

    for group in groups.values():
        if len(group) <= MAX_IMAGES_PER_SLIDE:
            result.extend(group)
            continue

        def _sz(it):
            try:
                from PIL import Image

                img = Image.open(io.BytesIO(it["bytes"]))
                return img.width * img.height

            except Exception:
                return 0

        result.extend(
            sorted(group, key=_sz, reverse=True)[:MAX_IMAGES_PER_SLIDE]
        )

    return result


def _filter_processed_images(
    processed_list: list[tuple],
    vision_mode: str,
) -> list[tuple]:
    """
    processed_list = [(ProcessedImage, normalized_bytes, original_item), ...]
    """
    from modules.extraction.visual.image import VisualType

    mode = _normalize_vision_mode(vision_mode)

    if mode == "text_only":
        return []

    priority = {
        VisualType.GRAPHIQUE: 100,
        VisualType.SCHEMA_TECHNIQUE: 95,
        VisualType.EQUATION: 90,
        VisualType.TABLEAU_IMAGE: 75,
        VisualType.CAPTURE_ECRAN: 55,
        VisualType.PLAN: 55,
        VisualType.PHOTO: 30,
        VisualType.INCONNU: 10,
    }

    kept = []

    for p, norm, item in processed_list:
        if p.skip_vision:
            continue

        vtype = p.visual_type
        conf = float(p.visual_type_confidence or 0.0)

        if mode == "auto":
            if vtype not in {
                VisualType.GRAPHIQUE,
                VisualType.SCHEMA_TECHNIQUE,
                VisualType.EQUATION,
                VisualType.TABLEAU_IMAGE,
            }:
                continue

            if vtype != VisualType.EQUATION and conf < 0.35:
                continue

        elif mode == "fast":
            if vtype not in {
                VisualType.EQUATION,
                VisualType.SCHEMA_TECHNIQUE,
                VisualType.GRAPHIQUE,
                VisualType.TABLEAU_IMAGE,
            }:
                continue

        elif mode == "full":
            pass

        kept.append((p, norm, item))

    if mode == "auto":
        limit = MAX_IMAGES_PER_DOCUMENT_AUTO
    elif mode == "fast":
        limit = MAX_IMAGES_PER_DOCUMENT_FAST
    else:
        limit = MAX_IMAGES_PER_DOCUMENT_FULL

    kept = sorted(
        kept,
        key=lambda x: (
            priority.get(x[0].visual_type, 0),
            float(x[0].visual_type_confidence or 0.0),
            x[0].width * x[0].height,
        ),
        reverse=True,
    )[:limit]

    return kept


def _inject_images(
    chunks: list[str],
    items: list[dict],
    source_type: str,
    vision_mode: str,
    formula_mode: str = "off",
) -> tuple[list[str], list[str]]:
    from modules.extraction.visual.image import process_image_bytes, VisualType
    from modules.extraction.visual.vision import describe_image_batch

    mode = _normalize_vision_mode(vision_mode)

    if not items or mode == "text_only":
        return chunks, []

    processed_list = []

    for item in items:
        norm = _safe_normalize_image(item["bytes"])

        if norm is None:
            continue

        try:
            p = process_image_bytes(
                norm,
                source_type,
                page_number=item.get("page"),
                slide_number=item.get("slide"),
                filename=item.get("filename"),
                caption=item.get("caption"),
            )

            processed_list.append((p, norm, item))

        except Exception as exc:
            logger.debug("Image ignorée : %s", exc)

    processed_list = _filter_processed_images(processed_list, mode)

    if not processed_list:
        return chunks, []

    try:
        from modules.extraction.visual.vision import warmup_vision_backend

        warmup_vision_backend()

    except Exception as exc:
        logger.debug("Warmup vision ignoré : %s", exc)

    cidx = _build_chunk_index(chunks)
    context_hints: dict[int, str] = {}

    for p, _, item in processed_list:
        num = item.get("slide") or item.get("page")
        tidx = cidx.get(num) if num else None

        if tidx is not None and 0 <= tidx < len(chunks):
            context_hints[id(p)] = chunks[tidx][:600]

    results = describe_image_batch(
        [p for p, _, _ in processed_list],
        context_hints=context_hints,
    )

    enriched = list(chunks)
    orphans: list[str] = []

    for (p, norm, item), vr in zip(processed_list, results):
        if not vr.description:
            continue

        num = item.get("slide") or item.get("page")
        tidx = cidx.get(num) if num else None

        fml_section = ""

        if (
            _normalize_formula_mode(formula_mode) != "off"
            and p.visual_type == VisualType.EQUATION
        ):
            try:
                from modules.extraction.formula.formula import extract_formulas

                fmls = extract_formulas(
                    image_bytes=norm,
                    page_number=item.get("page"),
                    formula_mode=formula_mode,
                )

                if fmls:
                    fml_section = "\n[FORMULES DÉTECTÉES]\n" + "".join(
                        f"  {i}. LaTeX: {f.latex}\n"
                        for i, f in enumerate(fmls, 1)
                    )

            except Exception as exc:
                logger.debug("Formule image ignorée : %s", exc)

        addition = f"\n[IMAGES]\n  • {vr.description}" + fml_section

        if tidx is not None:
            enriched[tidx] = enriched[tidx].rstrip() + addition
        else:
            orphans.append(vr.description)

    return enriched, orphans



# ──────────────────────────────────────────────────────────────────────────────
# Archives ZIP
# ──────────────────────────────────────────────────────────────────────────────

def _safe_archive_member_path(base_dir: Path, member_name: str) -> Path | None:
    """
    Empêche le Zip Slip : un membre ZIP ne doit jamais sortir du dossier temporaire.
    """
    raw = str(member_name or "").replace("\\", "/").strip()

    if not raw or raw.endswith("/"):
        return None

    # Ignore les fichiers système inutiles.
    lower = raw.lower()
    if lower.startswith("__macosx/") or lower.endswith(".ds_store"):
        return None

    candidate = (base_dir / raw).resolve()
    base_resolved = base_dir.resolve()

    try:
        candidate.relative_to(base_resolved)
    except Exception:
        return None

    return candidate


def _is_supported_archive_member(path: Path) -> bool:
    """
    Détermine si un fichier extrait du ZIP peut être repassé dans le router.
    """
    ext = path.suffix.lower()
    return (
        ext in EXTENSION_MAP
        or ext in AUDIO_VIDEO_EXTENSIONS
        or ext in ARCHIVE_EXTENSIONS
    )


def _run_archive(
    path: Path,
    source_tag: SourceTag,
    vision_mode: str,
    formula_mode: str = "off",
    enable_formulas: Optional[bool] = None,
    enable_transcription: bool = True,
    transcription_model: str = "small",
    transcription_language: Optional[str] = "fr",
    transcription_beam_size: int = 5,
    transcription_group_chunks: bool = True,
    transcription_chunk_seconds: int = DEFAULT_TRANSCRIPTION_CHUNK_SECONDS,
    transcription_chunk_max_chars: int = DEFAULT_TRANSCRIPTION_CHUNK_MAX_CHARS,
    archive_depth: int = 0,
    max_archive_depth: int = 2,
) -> ExtractionResult:
    """
    ZIP = conteneur documentaire :
    - extrait les fichiers supportés ;
    - repasse chaque fichier dans extract() ;
    - fusionne les chunks avec préfixe de traçabilité.

    Le NLP ne doit pas être modifié : il reçoit un ExtractionResult standard
    avec text_chunks / visual_chunks déjà enrichis.
    """
    category = _archive_file_category()

    result = ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=category,
        text_chunks=[],
        visual_chunks=[],
        detected_rd_sections=[],
        tags=["ARCHIVE:ZIP"],
        confidence_score=1.0,
        extraction_errors=[],
    )

    if archive_depth >= max_archive_depth:
        result.extraction_errors.append(
            f"Archive ignorée : profondeur maximale atteinte ({max_archive_depth})"
        )
        result.confidence_score = 0.2
        return result

    try:
        if not zipfile.is_zipfile(str(path)):
            result.extraction_errors.append("Archive ZIP invalide ou corrompue")
            result.confidence_score = 0.2
            return result

        tmp_dir = Path(tempfile.mkdtemp(prefix="ennosmart_zip_"))
        extracted_files: list[Path] = []
        total_size = 0

        with zipfile.ZipFile(str(path), "r") as zf:
            infos = [i for i in zf.infolist() if not i.is_dir()]

            if len(infos) > MAX_ARCHIVE_FILES:
                result.extraction_errors.append(
                    f"Archive limitée : {len(infos)} fichiers trouvés, maximum {MAX_ARCHIVE_FILES}"
                )
                infos = infos[:MAX_ARCHIVE_FILES]

            for info in infos:
                try:
                    member_path = _safe_archive_member_path(tmp_dir, info.filename)

                    if member_path is None:
                        continue

                    if info.file_size <= 0:
                        continue

                    if info.file_size > MAX_ARCHIVE_SINGLE_FILE_BYTES:
                        result.extraction_errors.append(
                            f"Fichier ZIP ignoré car trop volumineux : {info.filename} "
                            f"({round(info.file_size / (1024 * 1024), 2)} MB)"
                        )
                        continue

                    total_size += int(info.file_size or 0)

                    if total_size > MAX_ARCHIVE_TOTAL_BYTES:
                        result.extraction_errors.append(
                            "Archive limitée : taille totale maximale atteinte"
                        )
                        break

                    member_path.parent.mkdir(parents=True, exist_ok=True)

                    with zf.open(info, "r") as src, open(member_path, "wb") as dst:
                        dst.write(src.read())

                    if _is_supported_archive_member(member_path):
                        extracted_files.append(member_path)
                    else:
                        result.extraction_errors.append(
                            f"Fichier ZIP ignoré type non supporté : {info.filename}"
                        )

                except Exception as exc:
                    result.extraction_errors.append(
                        f"Erreur extraction membre ZIP {info.filename}: {exc}"
                    )

        logger.info(
            "Archive ZIP %s : %d fichiers supportés extraits",
            path.name,
            len(extracted_files),
        )

        if not extracted_files:
            result.extraction_errors.append("Archive ZIP : aucun fichier exploitable trouvé")
            result.confidence_score = 0.2
            return result

        chunks_added = 0

        for member_path in extracted_files:
            try:
                logger.info("Extraction fichier ZIP : %s", member_path.name)

                ar = extract(
                    member_path,
                    source_tag=source_tag,
                    vision_mode=vision_mode,
                    formula_mode=formula_mode,
                    enable_formulas=enable_formulas,
                    enable_transcription=enable_transcription,
                    transcription_model=transcription_model,
                    transcription_language=transcription_language,
                    transcription_beam_size=transcription_beam_size,
                    transcription_group_chunks=transcription_group_chunks,
                    transcription_chunk_seconds=transcription_chunk_seconds,
                    transcription_chunk_max_chars=transcription_chunk_max_chars,
                    archive_depth=archive_depth + 1,
                    max_archive_depth=max_archive_depth,
                )

                prefixed_text_chunks = [
                    f"[FICHIER ZIP : {member_path.name}]\n"
                    f"[ARCHIVE SOURCE : {path.name}]\n\n{chunk}"
                    for chunk in (ar.text_chunks or [])
                    if str(chunk or "").strip()
                ]

                prefixed_visual_chunks = [
                    f"[FICHIER ZIP : {member_path.name}]\n"
                    f"[ARCHIVE SOURCE : {path.name}]\n\n{chunk}"
                    for chunk in (ar.visual_chunks or [])
                    if str(chunk or "").strip()
                ]

                result.text_chunks.extend(prefixed_text_chunks)
                result.visual_chunks.extend(prefixed_visual_chunks)

                chunks_added += len(prefixed_text_chunks) + len(prefixed_visual_chunks)

                if getattr(ar, "detected_rd_sections", None):
                    result.detected_rd_sections.extend(ar.detected_rd_sections)

                result.tags.append(f"ZIP_FILE_EXTRACTED:{member_path.name}")

                if ar.tags:
                    for tag in ar.tags:
                        result.tags.append(f"ZIP:{member_path.name}:{tag}")

                if ar.extraction_errors:
                    for e in ar.extraction_errors:
                        # Si le fichier a quand même produit des chunks, certains messages
                        # sont de simples warnings techniques et ne doivent pas polluer ERRORS.
                        if (ar.text_chunks or ar.visual_chunks) and _is_benign_extraction_warning(e):
                            result.tags.append(f"ZIP_WARNING:{member_path.name}:{e}")
                            continue

                        result.extraction_errors.append(f"ZIP {member_path.name}: {e}")

                if not prefixed_text_chunks and not prefixed_visual_chunks:
                    result.extraction_errors.append(
                        f"ZIP {member_path.name}: aucun chunk exploitable extrait"
                    )

            except Exception as exc:
                logger.warning("Erreur extraction fichier ZIP %s : %s", member_path, exc)
                result.extraction_errors.append(f"ZIP {member_path.name}: {exc}")

        if chunks_added:
            result.tags.append("ARCHIVE_FILES_EXTRACTED")
            result.tags.append(f"ARCHIVE_FILE_CHUNKS:{chunks_added}")
            result.confidence_score = 1.0
        else:
            result.confidence_score = 0.2

        result.tags = list(dict.fromkeys(result.tags or []))
        result.detected_rd_sections = list(dict.fromkeys(result.detected_rd_sections or []))

        return result

    except Exception as exc:
        logger.error("Erreur archive ZIP %s : %s", path.name, exc, exc_info=True)
        result.extraction_errors.append(str(exc))
        result.confidence_score = 0.2
        return result


# ──────────────────────────────────────────────────────────────────────────────
# Pipelines par type fichier
# ──────────────────────────────────────────────────────────────────────────────

def _run_pdf(
    path: Path,
    vision_mode: str,
    formula_mode: str = "off",
) -> ExtractionResult:
    from modules.extraction.text.pdf_native import extract_pdf_native
    from modules.extraction.text.pdf_ocr import extract_pdf_ocr, merge_native_and_ocr

    native = extract_pdf_native(str(path))
    final_chunks = native.text_chunks
    confidence = native.confidence_score

    if native.ocr_needed_pages:
        ocr = extract_pdf_ocr(
            str(path),
            target_pages=native.ocr_needed_pages,
        )

        final_chunks = merge_native_and_ocr(
            native.text_chunks,
            ocr,
        )

        native.extraction_errors.extend(ocr.extraction_errors)
        confidence = native.confidence_score * 0.7 + ocr.confidence_score * 0.3

    enriched = _inject_formulas(
        final_chunks,
        path.name,
        formula_mode=formula_mode,
    )

    image_items = []

    if _normalize_vision_mode(vision_mode) != "text_only":
        try:
            import fitz

            doc = fitz.open(str(path))

            for pr in [p for p in native.pages if getattr(p, "has_images", False)]:
                for img_info in doc[pr.page_number - 1].get_images(full=True):
                    image_items.append(
                        {
                            "bytes": doc.extract_image(img_info[0])["image"],
                            "page": pr.page_number,
                            "slide": None,
                        }
                    )

            doc.close()

        except Exception as exc:
            logger.warning("Images PDF : %s", exc)

    image_items = _limit_images_per_slide(image_items)

    enriched, orphans = _inject_images(
        enriched,
        image_items,
        "pdf_page",
        vision_mode,
        formula_mode=formula_mode,
    )

    tags = list(set(native.tags))

    if _normalize_formula_mode(formula_mode) == "off":
        tags.append("FORMULAS_DISABLED")
    else:
        tags.append(f"FORMULA_MODE:{formula_mode.upper()}")

    if _normalize_vision_mode(vision_mode) == "text_only":
        tags.append("VISUALS_DISABLED")
    else:
        tags.append(f"VISION_MODE:{_normalize_vision_mode(vision_mode).upper()}")

    if native.ocr_needed_pages:
        tags.append("MIXED_NATIVE_OCR")

    if orphans or any("[IMAGES]" in c for c in enriched):
        tags.append("HAS_VISUAL_DESCRIPTIONS")

    if any("[FORMULES" in c for c in enriched):
        tags.append("HAS_FORMULAS")

    return ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=(
            FileCategory.PDF_NATIVE
            if not native.ocr_needed_pages
            else FileCategory.PDF_OCR
        ),
        text_chunks=enriched,
        visual_chunks=orphans,
        title=native.metadata.title,
        author=native.metadata.author,
        creation_date=native.metadata.creation_date,
        page_count=native.metadata.page_count,
        detected_rd_sections=native.detected_rd_sections,
        tags=tags,
        confidence_score=round(confidence, 2),
        extraction_errors=native.extraction_errors,
        ocr_needed_pages=native.ocr_needed_pages,
    )


def _run_office(
    path: Path,
    vision_mode: str,
    formula_mode: str = "off",
) -> ExtractionResult:
    from modules.extraction.text.office import extract_office

    office = extract_office(str(path))

    enriched = _inject_formulas(
        office.text_chunks,
        path.name,
        formula_mode=formula_mode,
    )

    if _normalize_formula_mode(formula_mode) != "off" and office.file_type in {"docx", "docm"}:
        enriched = _run_docx_omml(
            path,
            enriched,
            formula_mode=formula_mode,
        )

    elif _normalize_formula_mode(formula_mode) != "off" and office.file_type in {"pptx", "pptm"}:
        enriched = _run_pptx_omml(
            path,
            enriched,
            formula_mode=formula_mode,
        )

    image_items = []

    if _normalize_vision_mode(vision_mode) != "text_only":
        try:
            if office.file_type in {"docx", "docm"}:
                import zipfile

                with zipfile.ZipFile(str(path)) as zf:
                    for mf in [
                        f
                        for f in zf.namelist()
                        if f.startswith("word/media/")
                    ]:
                        image_items.append(
                            {
                                "bytes": zf.read(mf),
                                "page": None,
                                "slide": None,
                                "filename": Path(mf).name,
                            }
                        )

            elif office.visual_candidates:
                from pptx import Presentation

                prs = Presentation(str(path))

                for snum in office.visual_candidates:
                    for shape in prs.slides[snum - 1].shapes:
                        if hasattr(shape, "image"):
                            image_items.append(
                                {
                                    "bytes": shape.image.blob,
                                    "slide": snum,
                                    "page": None,
                                }
                            )

        except Exception as exc:
            logger.warning("Collecte images office : %s", exc)

    image_items = _limit_images_per_slide(image_items)

    enriched, orphans = _inject_images(
        enriched,
        image_items,
        "docx_document" if office.file_type in {"docx", "docm"} else "pptx_slide",
        vision_mode,
        formula_mode=formula_mode,
    )

    tags = list(set(office.tags))

    if _normalize_formula_mode(formula_mode) == "off":
        tags.append("FORMULAS_DISABLED")
    else:
        tags.append(f"FORMULA_MODE:{formula_mode.upper()}")

    if _normalize_vision_mode(vision_mode) == "text_only":
        tags.append("VISUALS_DISABLED")
    else:
        tags.append(f"VISION_MODE:{_normalize_vision_mode(vision_mode).upper()}")

    if orphans or any("[IMAGES]" in c for c in enriched):
        tags.append("HAS_VISUAL_DESCRIPTIONS")

    if any("[FORMULES" in c for c in enriched):
        tags.append("HAS_FORMULAS")

    return ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=(
            FileCategory.DOCX
            if office.file_type in {"docx", "docm"}
            else FileCategory.PPTX
        ),
        text_chunks=enriched,
        visual_chunks=orphans,
        title=office.metadata.title,
        author=office.metadata.author,
        page_count=office.metadata.slide_count or 0,
        detected_rd_sections=office.detected_rd_sections,
        tags=tags,
        confidence_score=office.confidence_score,
        extraction_errors=office.extraction_errors,
    )



def _safe_temp_attachment_name(filename: str, fallback: str = "attachment.bin") -> str:
    """
    Nettoie le nom d'une pièce jointe avant écriture temporaire.
    Le router ne doit jamais faire confiance au chemin fourni par le mail.
    """
    name = str(filename or "").replace("\\", "/").split("/")[-1].strip()
    if not name:
        name = fallback

    name = re.sub(r"[\x00-\x1f\x7f]+", "_", name)
    name = re.sub(r"[<>:\"/\\|?*]+", "_", name)
    name = re.sub(r"\s+", " ", name).strip()

    return name or fallback



def _is_benign_extraction_warning(message: str) -> bool:
    """
    Messages qui ne doivent pas faire échouer un document.
    Exemple : python-docx refuse un .docm, mais office.py extrait ensuite le contenu via XML fallback.
    """
    msg = str(message or "").lower()

    if "fallback xml utilisé" in msg or "xml fallback" in msg:
        if "python-docx" in msg or "macroenabled" in msg:
            return True

    if "aucun chunk exploitable extrait" in msg:
        # Ce message peut être bénin pour une image/logo en mode text_only.
        return True

    return False


def _is_image_attachment_name(filename: str) -> bool:
    ext = Path(str(filename or "")).suffix.lower()
    return ext in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp", ".svg"}


def _run_email(
    path: Path,
    formula_mode: str = "off",
) -> ExtractionResult:
    """
    Email (.msg/.eml) = conteneur :
    - extrait le corps du mail ;
    - conserve les pièces jointes R&D en fichiers temporaires ;
    - le point d'entrée extract() les repasse ensuite dans le routeur.

    Important :
    Le parser email ne doit pas analyser les pièces jointes lui-même.
    Il fournit seulement attachments_paths au router.
    """
    from modules.extraction.text.email_parser import extract_email

    email_result = extract_email(str(path))
    attachment_paths: list[str] = []
    attachment_errors: list[str] = []

    rd_attachments = [
        a for a in email_result.attachments
        if getattr(a, "is_rd_relevant", False)
    ]

    logger.info(
        "Email %s : %d pièces jointes détectées, %d pertinentes R&D",
        path.name,
        len(email_result.attachments),
        len(rd_attachments),
    )

    for idx, att in enumerate(rd_attachments, start=1):
        try:
            content = getattr(att, "content", None)

            if not content:
                msg = (
                    f"Pièce jointe sans contenu exploitable : "
                    f"{getattr(att, 'filename', 'unknown')}"
                )
                logger.warning(msg)
                attachment_errors.append(msg)
                continue

            safe_name = _safe_temp_attachment_name(
                getattr(att, "filename", "") or f"attachment_{idx}.bin",
                fallback=f"attachment_{idx}.bin",
            )

            tmp_dir = Path(tempfile.mkdtemp(prefix="ennosmart_att_"))
            tmp_path = tmp_dir / safe_name
            tmp_path.write_bytes(content)

            attachment_paths.append(str(tmp_path))

            logger.info(
                "PJ email prête pour extraction : %s | ext=%s | size=%d | tmp=%s",
                safe_name,
                getattr(att, "extension", Path(safe_name).suffix.lower()),
                len(content),
                tmp_path,
            )

        except Exception as exc:
            msg = f"Erreur préparation PJ email {getattr(att, 'filename', 'unknown')} : {exc}"
            logger.warning(msg)
            attachment_errors.append(msg)

    enriched = _inject_formulas(
        email_result.text_chunks,
        path.name,
        formula_mode=formula_mode,
    )

    tags = list(set(email_result.tags))

    if attachment_paths:
        tags.append("EMAIL_ATTACHMENTS_READY")
        tags.append(f"EMAIL_ATTACHMENTS_READY_COUNT:{len(attachment_paths)}")

    if rd_attachments:
        tags.append(f"EMAIL_RD_ATTACHMENTS_COUNT:{len(rd_attachments)}")

    if _normalize_formula_mode(formula_mode) == "off":
        tags.append("FORMULAS_DISABLED")
    else:
        tags.append(f"FORMULA_MODE:{formula_mode.upper()}")

    if any("[FORMULES" in c for c in enriched):
        tags.append("HAS_FORMULAS")

    extraction_errors = list(email_result.extraction_errors or []) + attachment_errors

    return ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=FileCategory.EMAIL,
        text_chunks=enriched,
        attachments_paths=attachment_paths,
        detected_rd_sections=email_result.detected_rd_sections,
        tags=list(dict.fromkeys(tags)),
        confidence_score=email_result.confidence_score,
        extraction_errors=extraction_errors,
    )


def _run_excel(
    path: Path,
    formula_mode: str = "off",
) -> ExtractionResult:
    from modules.extraction.structured.excel_struct import extract_excel

    excel = extract_excel(str(path))
    enriched = list(excel.text_chunks)

    if _normalize_formula_mode(formula_mode) != "off":
        try:
            from modules.extraction.formula.formula import extract_formulas

            for sheet in excel.sheets:
                sidx = next(
                    (
                        i
                        for i, c in enumerate(enriched)
                        if sheet.name in c
                    ),
                    len(enriched) - 1,
                )

                seen: set[str] = set()
                section = ""

                for addr, cell in sheet.raw_cells.items():
                    value = str(cell.value or "")

                    if value.startswith("="):
                        for r in extract_formulas(
                            excel_formula=value,
                            formula_mode=formula_mode,
                        ):
                            k = re.sub(r"\s+", "", r.latex.strip().lower())

                            if k and k not in seen:
                                seen.add(k)
                                section += f"  • {addr}: {r.latex}\n"

                if section and 0 <= sidx < len(enriched):
                    enriched[sidx] = (
                        enriched[sidx].rstrip()
                        + "\n[FORMULES DÉTECTÉES]\n"
                        + section
                    )

        except Exception as exc:
            logger.debug("Formules Excel ignorées : %s", exc)

    tags = list(set(excel.tags))

    if _normalize_formula_mode(formula_mode) == "off":
        tags.append("FORMULAS_DISABLED")
    else:
        tags.append(f"FORMULA_MODE:{formula_mode.upper()}")

    if any("[FORMULES" in c for c in enriched):
        tags.append("HAS_FORMULAS")

    return ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=FileCategory.EXCEL,
        text_chunks=enriched,
        structured_data=excel.structured_data,
        detected_rd_sections=excel.detected_rd_sections,
        tags=tags,
        confidence_score=excel.confidence_score,
        extraction_errors=excel.extraction_errors,
    )


def _run_image(
    path: Path,
    vision_mode: str,
    formula_mode: str = "off",
) -> ExtractionResult:
    from modules.extraction.visual.image import process_image_file, VisualType
    from modules.extraction.visual.vision import describe_image

    processed = process_image_file(str(path))

    if _normalize_vision_mode(vision_mode) == "text_only":
        return ExtractionResult(
            file_name=path.name,
            source_path=str(path.resolve()),
            file_category=FileCategory.IMAGE,
            text_chunks=[],
            visual_chunks=[],
            tags=list(set(processed.tags + ["VISUAL_DISABLED"])),
            confidence_score=0.0,
            extraction_errors=processed.processing_errors,
        )

    vr = describe_image(processed)
    chunk = vr.description or ""

    if (
        _normalize_formula_mode(formula_mode) != "off"
        and processed.visual_type == VisualType.EQUATION
        and chunk
    ):
        try:
            from modules.extraction.formula.formula import extract_formulas

            fmls = extract_formulas(
                image_bytes=processed.image_bytes,
                formula_mode=formula_mode,
            )

            if fmls:
                chunk = (
                    chunk.rstrip()
                    + "\n[FORMULES DÉTECTÉES]\n"
                    + "".join(
                        f"  {i}. LaTeX: {f.latex}\n"
                        for i, f in enumerate(fmls, 1)
                    )
                )

        except Exception as exc:
            logger.debug("Formule image ignorée : %s", exc)

    tags = list(set(processed.tags + vr.tags))

    if _normalize_formula_mode(formula_mode) == "off":
        tags.append("FORMULAS_DISABLED")
    else:
        tags.append(f"FORMULA_MODE:{formula_mode.upper()}")

    if "[FORMULES" in chunk:
        tags.append("HAS_FORMULAS")

    return ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=FileCategory.IMAGE,
        text_chunks=[],
        visual_chunks=[chunk] if chunk else [],
        tags=tags,
        confidence_score=1.0 if not processed.skip_vision else 0.0,
        extraction_errors=processed.processing_errors + ([vr.error] if vr.error else []),
    )


# ──────────────────────────────────────────────────────────────────────────────
# Audio / Vidéo — transcription
# ──────────────────────────────────────────────────────────────────────────────

def _format_time(seconds: float) -> str:
    seconds = max(0.0, float(seconds or 0.0))
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    return f"{h:02d}:{m:02d}:{s:02d}"


def _engine_value(engine) -> str:
    if engine is None:
        return "faster_whisper"

    return getattr(engine, "value", str(engine))


def _build_grouped_transcription_chunks(
    path: Path,
    tr,
    max_seconds: int = DEFAULT_TRANSCRIPTION_CHUNK_SECONDS,
    max_chars: int = DEFAULT_TRANSCRIPTION_CHUNK_MAX_CHARS,
) -> list[str]:
    """
    Regroupe les segments Whisper en chunks RAG plus stables.

    Exemple :
    - au lieu de 963 segments = 963 chunks
    - on obtient environ 35 à 60 chunks pour un audio d'une heure.
    """
    segments = list(getattr(tr, "segments", []) or [])

    if not segments:
        return list(getattr(tr, "text_chunks", []) or [])

    engine = _engine_value(getattr(tr, "engine_used", None))
    chunks: list[str] = []

    current_texts: list[str] = []
    current_start: Optional[float] = None
    current_end: Optional[float] = None
    current_chars = 0
    chunk_index = 1

    def flush() -> None:
        nonlocal current_texts, current_start, current_end, current_chars, chunk_index

        text = " ".join(t.strip() for t in current_texts if t and t.strip()).strip()

        if not text:
            current_texts = []
            current_start = None
            current_end = None
            current_chars = 0
            return

        start = 0.0 if current_start is None else current_start
        end = start if current_end is None else current_end

        chunk = (
            f"[AUDIO:{path.name}] "
            f"[CHUNK {chunk_index}] "
            f"[{_format_time(start)} → {_format_time(end)}] "
            f"[TRANSCRIPTION:{engine}]\n\n"
            f"{text}"
        )

        chunks.append(chunk)

        chunk_index += 1
        current_texts = []
        current_start = None
        current_end = None
        current_chars = 0

    for seg in segments:
        text = str(getattr(seg, "text", "") or "").strip()

        if not text:
            continue

        start = float(getattr(seg, "start", 0.0) or 0.0)
        end = float(getattr(seg, "end", start) or start)

        if current_start is None:
            current_start = start

        should_flush = False

        if current_end is not None:
            duration_if_added = end - current_start
            chars_if_added = current_chars + len(text)

            if duration_if_added >= max_seconds:
                should_flush = True

            if chars_if_added >= max_chars:
                should_flush = True

        if should_flush:
            flush()
            current_start = start

        current_texts.append(text)
        current_end = end
        current_chars += len(text) + 1

    flush()

    return chunks


def _run_audio(
    path: Path,
    transcription_model: str = "small",
    transcription_language: Optional[str] = "fr",
    transcription_beam_size: int = 5,
    transcription_group_chunks: bool = True,
    transcription_chunk_seconds: int = DEFAULT_TRANSCRIPTION_CHUNK_SECONDS,
    transcription_chunk_max_chars: int = DEFAULT_TRANSCRIPTION_CHUNK_MAX_CHARS,

    # Archives ZIP
    archive_depth: int = 0,
    max_archive_depth: int = 2,
) -> ExtractionResult:
    """
    Audio/vidéo → transcription texte.

    Résultat aligné avec les autres extracteurs :
    - text_chunks : chunks horodatés utilisables par NLP + RAG
    - visual_chunks : vide
    - tags : transcription, modèle, langue, device
    - métadonnées transcription remplies dans ExtractionResult
    """
    try:
        from modules.extraction.audio.audio_transcriber import extract_audio_transcription
    except Exception as exc:
        logger.error("Module transcription indisponible : %s", exc)

        return ExtractionResult(
            file_name=path.name,
            source_path=str(path.resolve()),
            file_category=_audio_file_category(),
            text_chunks=[],
            visual_chunks=[],
            tags=[
                "TRANSCRIPTION:NONE",
                "ERROR:TRANSCRIBER_IMPORT_FAILED",
            ],
            confidence_score=0.0,
            extraction_errors=[
                "Module audio_transcriber indisponible. "
                "Vérifier modules/extraction/audio/audio_transcriber.py "
                "et installer faster-whisper.",
                str(exc),
            ],
        )

    tr = extract_audio_transcription(
        file_path=str(path),
        model_name=transcription_model,
        language=transcription_language,
        beam_size=transcription_beam_size,
    )

    if transcription_group_chunks:
        chunks = _build_grouped_transcription_chunks(
            path=path,
            tr=tr,
            max_seconds=transcription_chunk_seconds,
            max_chars=transcription_chunk_max_chars,
        )
    else:
        chunks = list(getattr(tr, "text_chunks", []) or [])

    if not chunks and getattr(tr, "full_text", ""):
        chunks = [
            (
                f"[AUDIO:{path.name}] "
                f"[TRANSCRIPTION:{_engine_value(getattr(tr, 'engine_used', None))}]\n\n"
                f"{tr.full_text.strip()}"
            )
        ]

    duration = getattr(tr, "duration", None)
    language = getattr(tr, "language", None)
    model_name = getattr(tr, "model_name", None)
    engine = _engine_value(getattr(tr, "engine_used", None))

    tags = list(set(getattr(tr, "tags", []) or []))
    tags.append("AUDIO_VIDEO_TRANSCRIPTION")

    if _is_video(path):
        tags.append("SOURCE:VIDEO")
    else:
        tags.append("SOURCE:AUDIO")

    if language:
        tags.append(f"TRANSCRIPTION_LANG:{language}")

    if model_name:
        tags.append(f"TRANSCRIPTION_MODEL:{model_name}")

    if duration is not None:
        try:
            tags.append(f"DURATION_SECONDS:{round(float(duration), 2)}")
        except Exception:
            pass

    if transcription_group_chunks:
        tags.append("TRANSCRIPTION_CHUNKS_GROUPED")
        tags.append(f"TRANSCRIPTION_CHUNK_SECONDS:{transcription_chunk_seconds}")

    errors = list(getattr(tr, "extraction_errors", []) or [])

    if not chunks:
        errors.append("Transcription vide : aucun texte exploitable extrait")

    return ExtractionResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_category=_audio_file_category(),
        text_chunks=chunks,
        visual_chunks=[],
        tags=list(set(tags)),
        confidence_score=float(getattr(tr, "confidence_score", 0.0) or 0.0),
        extraction_errors=errors,

        # Métadonnées audio/vidéo corrigées
        media_duration_seconds=(
            float(duration)
            if duration is not None
            else None
        ),
        transcription_language=language,
        transcription_model=model_name,
        transcription_engine=engine,
    )


# ──────────────────────────────────────────────────────────────────────────────
# Point d'entrée
# ──────────────────────────────────────────────────────────────────────────────

def extract(
    file_path,
    source_tag: SourceTag = SourceTag.DE_DOC,
    vision_mode: str = "text_only",
    formula_mode: str = "off",
    enable_formulas: Optional[bool] = None,
    mode: Optional[str] = None,

    # Transcription audio/vidéo
    enable_transcription: bool = True,
    transcription_model: str = "small",
    transcription_language: Optional[str] = "fr",
    transcription_beam_size: int = 5,
    transcription_group_chunks: bool = True,
    transcription_chunk_seconds: int = DEFAULT_TRANSCRIPTION_CHUNK_SECONDS,
    transcription_chunk_max_chars: int = DEFAULT_TRANSCRIPTION_CHUNK_MAX_CHARS,

    # Archives ZIP
    archive_depth: int = 0,
    max_archive_depth: int = 2,
) -> ExtractionResult:
    """
    Paramètres recommandés :
      vision_mode  : text_only | auto | fast | full
      formula_mode : off | fast | explain

    Défaut evidence-first / tests NLP :
      vision_mode="text_only"
      formula_mode="off"

    Pour production enrichie :
      vision_mode="full"
      formula_mode="fast"

    Audio/vidéo :
      enable_transcription=True
      transcription_model="small"
      transcription_language="fr"
      transcription_group_chunks=True

    Compatibilité :
      - enable_formulas=False équivaut à formula_mode="off"
      - mode="text" ou "text_only" désactive vision
      - mode="fast" garde vision fast si vision_mode non fourni
      - mode="full" garde vision full si demandé explicitement par CLI
    """
    path = Path(file_path)

    # Compatibilité avec anciens scripts : --mode text/fast/full
    if mode:
        m = str(mode).lower().strip()

        if m in {"text", "text_only"}:
            vision_mode = "text_only"

        elif m == "fast" and vision_mode in {"auto", "", None}:
            vision_mode = "fast"

        elif m == "full" and vision_mode in {"", None}:
            vision_mode = "full"

    vision_mode = _normalize_vision_mode(vision_mode)
    formula_mode = _normalize_formula_mode(
        formula_mode,
        enable_formulas=enable_formulas,
    )

    if not path.exists():
        return ExtractionResult(
            file_name=path.name,
            source_path=str(path),
            file_category=FileCategory.UNKNOWN,
            extraction_errors=[f"Fichier introuvable : {path}"],
        )

    category = _detect_category(path)

    logger.info(
        "Extraction : %s [%s] vision=%s formula=%s transcription=%s",
        path.name,
        _category_label(category, path),
        vision_mode,
        formula_mode,
        enable_transcription,
    )

    try:
        if path.suffix.lower() in ARCHIVE_EXTENSIONS:
            result = _run_archive(
                path=path,
                source_tag=source_tag,
                vision_mode=vision_mode,
                formula_mode=formula_mode,
                enable_formulas=enable_formulas,
                enable_transcription=enable_transcription,
                transcription_model=transcription_model,
                transcription_language=transcription_language,
                transcription_beam_size=transcription_beam_size,
                transcription_group_chunks=transcription_group_chunks,
                transcription_chunk_seconds=transcription_chunk_seconds,
                transcription_chunk_max_chars=transcription_chunk_max_chars,
                archive_depth=archive_depth,
                max_archive_depth=max_archive_depth,
            )

        elif _is_audio_video(path):
            if not enable_transcription:
                return ExtractionResult(
                    file_name=path.name,
                    source_path=str(path.resolve()),
                    file_category=_audio_file_category(),
                    extraction_errors=[
                        "Fichier audio/vidéo détecté mais transcription désactivée"
                    ],
                    tags=["TRANSCRIPTION_DISABLED"],
                )

            result = _run_audio(
                path,
                transcription_model=transcription_model,
                transcription_language=transcription_language,
                transcription_beam_size=transcription_beam_size,
                transcription_group_chunks=transcription_group_chunks,
                transcription_chunk_seconds=transcription_chunk_seconds,
                transcription_chunk_max_chars=transcription_chunk_max_chars,
            )

        elif category in (FileCategory.PDF_NATIVE, FileCategory.PDF_OCR):
            result = _run_pdf(
                path,
                vision_mode,
                formula_mode=formula_mode,
            )

        elif category in (FileCategory.DOCX, FileCategory.PPTX):
            result = _run_office(
                path,
                vision_mode,
                formula_mode=formula_mode,
            )

        elif category == FileCategory.EMAIL:
            result = _run_email(
                path,
                formula_mode=formula_mode,
            )

        elif category == FileCategory.EXCEL:
            result = _run_excel(
                path,
                formula_mode=formula_mode,
            )

        elif category == FileCategory.IMAGE:
            result = _run_image(
                path,
                vision_mode,
                formula_mode=formula_mode,
            )

        else:
            return ExtractionResult(
                file_name=path.name,
                source_path=str(path),
                file_category=FileCategory.UNKNOWN,
                extraction_errors=[f"Type non supporté : {path.suffix}"],
            )

        result.source_tag = source_tag

        # Pièces jointes email : extraction récursive par le même routeur.
        # Le mail est un conteneur ; ses PJ deviennent du contenu RAG/NLP.
        if result.attachments_paths:
            logger.info(
                "Extraction des pièces jointes email pour %s : %d PJ",
                path.name,
                len(result.attachments_paths),
            )

            attachment_chunks_added = 0

            for att_path in result.attachments_paths:
                att_name = Path(att_path).name

                try:
                    logger.info("Extraction PJ email : %s", att_name)

                    ar = extract(
                        att_path,
                        source_tag=source_tag,
                        vision_mode=vision_mode,
                        formula_mode=formula_mode,
                        enable_formulas=enable_formulas,
                        enable_transcription=enable_transcription,
                        transcription_model=transcription_model,
                        transcription_language=transcription_language,
                        transcription_beam_size=transcription_beam_size,
                        transcription_group_chunks=transcription_group_chunks,
                        transcription_chunk_seconds=transcription_chunk_seconds,
                        transcription_chunk_max_chars=transcription_chunk_max_chars,
                        archive_depth=archive_depth,
                        max_archive_depth=max_archive_depth,
                    )

                    # Préfixe de traçabilité : on garde l'origine PJ dans le texte.
                    prefixed_text_chunks = [
                        f"[PIÈCE JOINTE EMAIL : {att_name}]\n"
                        f"[EMAIL SOURCE : {path.name}]\n\n{chunk}"
                        for chunk in (ar.text_chunks or [])
                        if str(chunk or "").strip()
                    ]

                    prefixed_visual_chunks = [
                        f"[PIÈCE JOINTE EMAIL : {att_name}]\n"
                        f"[EMAIL SOURCE : {path.name}]\n\n{chunk}"
                        for chunk in (ar.visual_chunks or [])
                        if str(chunk or "").strip()
                    ]

                    result.text_chunks.extend(prefixed_text_chunks)
                    result.visual_chunks.extend(prefixed_visual_chunks)

                    attachment_chunks_added += len(prefixed_text_chunks) + len(prefixed_visual_chunks)

                    if getattr(ar, "detected_rd_sections", None):
                        result.detected_rd_sections.extend(ar.detected_rd_sections)

                    result.tags.append(f"EMAIL_ATTACHMENT_EXTRACTED:{att_name}")

                    if ar.tags:
                        for tag in ar.tags:
                            result.tags.append(f"PJ:{att_name}:{tag}")

                    if ar.extraction_errors:
                        for e in ar.extraction_errors:
                            if (ar.text_chunks or ar.visual_chunks) and _is_benign_extraction_warning(e):
                                result.tags.append(f"PJ_WARNING:{att_name}:{e}")
                                continue

                            result.extraction_errors.append(f"PJ {att_name}: {e}")

                    if not prefixed_text_chunks and not prefixed_visual_chunks:
                        # Une image/logo/signature en vision_mode=text_only ne doit pas compter
                        # comme erreur d'extraction. Elle est volontairement ignorée.
                        if _is_image_attachment_name(att_name) and _normalize_vision_mode(vision_mode) == "text_only":
                            result.tags.append(f"EMAIL_ATTACHMENT_SKIPPED_NO_TEXT:{att_name}")
                        else:
                            result.extraction_errors.append(
                                f"PJ {att_name}: aucun chunk exploitable extrait"
                            )

                except Exception as exc:
                    logger.warning("PJ %s : %s", att_path, exc)
                    result.extraction_errors.append(f"PJ {att_name}: {exc}")

            if attachment_chunks_added:
                result.tags.append("EMAIL_ATTACHMENTS_EXTRACTED")
                result.tags.append(f"EMAIL_ATTACHMENT_CHUNKS:{attachment_chunks_added}")

            # Déduplication simple des tags et sections.
            result.tags = list(dict.fromkeys(result.tags))
            result.detected_rd_sections = list(dict.fromkeys(result.detected_rd_sections or []))

        result.tags = list(dict.fromkeys(result.tags or []))
        result.detected_rd_sections = list(dict.fromkeys(result.detected_rd_sections or []))

        logger.info(
            "Résultat : %d chunks texte, %d visuels | tags=%s",
            len(result.text_chunks),
            len(result.visual_chunks),
            result.tags,
        )

        return result

    except Exception as exc:
        logger.error("Erreur extraction %s : %s", path.name, exc, exc_info=True)

        return ExtractionResult(
            file_name=path.name,
            source_path=str(path),
            file_category=category,
            extraction_errors=[str(exc)],
        )


# Alias compatibilité
extract_file = extract
extract_document = extract
process_file = extract
process_document = extract
run_extraction = extract
route_file = extract
extract_any = extract