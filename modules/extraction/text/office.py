"""
modules/extraction/text/office.py
──────────────────────────────────────────────────────────────────────────────
Extraction de documents Office (PPTX + DOCX) pour dossiers R&D / CIR.

Stratégie DOCX :
  - Extraction paragraphe par paragraphe avec préservation des styles
    (Titre, Heading1..6, Normal, ListParagraph)
  - Extraction des tableaux Word → Markdown (mêmes règles que pdf_native)
  - Détection des zones de commentaires / annotations (notes consultant)
  - Extraction des propriétés document (auteur, société, date révision)

Stratégie PPTX :
  - Extraction slide par slide → un chunk par slide
  - Texte des zones de texte + notes présentateur (souvent riches en R&D)
  - Extraction des titres de slides pour reconstituer le plan
  - Tableaux dans les slides → Markdown
  - Détection d'images/schémas → flag pour visual/vision.py
"""

from __future__ import annotations

import logging
import re
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Optional

try:
    from docx import Document as DocxDocument
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


# ── Constantes ────────────────────────────────────────────────────────────────

# Styles Word considérés comme titres/sections
DOCX_HEADING_STYLES = {
    "heading 1", "heading 2", "heading 3",
    "heading 4", "heading 5", "heading 6",
    "titre", "titre 1", "titre 2", "titre 3",   # Variantes FR
}

# Styles Word de contenu normal
DOCX_BODY_STYLES = {
    "normal", "default", "corps de texte",
    "list paragraph", "liste à puces",
    "no spacing", "body text",
}

# Seuil texte minimum par slide / paragraphe
MIN_CHARS_SLIDE    = 20
MIN_CHARS_PARA     = 5

# Sections R&D standards (mêmes patterns que pdf_native)
RD_SECTION_PATTERNS: list[str] = [
    r"(objectif[s]?\s+(?:du\s+)?projet)",
    r"(état\s+de\s+l['\s]art)",
    r"(verrous?\s+technologique[s]?)",
    r"(travaux\s+(?:de\s+)?r(?:echerche)?(?:\s*&\s*|\s+et\s+)d(?:éveloppement)?)",
    r"(résultats?\s+(?:obtenus?|attendus?))",
    r"(description\s+(?:des\s+)?travaux)",
    r"(dépenses?\s+(?:de\s+)?recherche)",
    r"(personnel\s+(?:de\s+)?recherche)",
    r"(nature\s+(?:de\s+)?l['\s]incertitude)",
    r"(démarche\s+(?:scientifique|expérimentale))",
]

_RD_SECTION_RE = re.compile(
    "|".join(RD_SECTION_PATTERNS),
    flags=re.IGNORECASE | re.UNICODE,
)


# ── Enums ─────────────────────────────────────────────────────────────────────

class OfficeFileType(str, Enum):
    DOCX = "docx"
    DOCM = "docm"
    PPTX = "pptx"
    PPTM = "pptm"


# ── Dataclasses ───────────────────────────────────────────────────────────────

@dataclass
class DocxBlockResult:
    """Un bloc extrait d'un DOCX (paragraphe ou tableau)."""
    block_index: int
    block_type: str            # "heading" | "paragraph" | "table" | "comment"
    style_name: str
    raw_text: str
    level: int                 # Niveau de titre (1-6), 0 si corps
    is_rd_section: bool
    char_count: int


@dataclass
class SlideResult:
    """Résultat d'extraction pour une slide PPTX."""
    slide_number: int
    title: Optional[str]
    raw_text: str              # Tout le texte des zones texte
    notes_text: str            # Notes présentateur
    tables_as_text: list[str]  # Tableaux → Markdown
    has_images: bool           # Images/schémas détectés
    has_charts: bool           # Graphiques détectés
    images_descriptions: list[dict] = field(default_factory=list)  # [{"index": 1, "description": "..."}]
    detected_sections: list[str] = field(default_factory=list)
    char_count: int = 0
    is_text_poor: bool = False


@dataclass
class DocumentMetadata:
    """Métadonnées communes DOCX / PPTX."""
    title: Optional[str]            = None
    author: Optional[str]           = None
    company: Optional[str]          = None
    last_modified_by: Optional[str] = None
    created: Optional[str]          = None
    modified: Optional[str]         = None
    revision: Optional[str]         = None
    organisme_detected: Optional[str] = None
    slide_count: int                = 0    # PPTX seulement
    word_count: int                 = 0    # DOCX seulement
    has_comments: bool              = False


@dataclass
class OfficeResult:
    """
    Résultat complet d'extraction d'un fichier Office.
    Compatible avec ExtractionResult (base.py).
    """
    file_name: str
    source_path: str
    file_type: str              # "docx" | "docm" | "pptx" | "pptm"

    # ── Sortie principale pour le RAG ──────────────────────────────────────
    text_chunks: list[str] = field(default_factory=list)
    # DOCX : 1 chunk par section logique (heading + contenu)
    # PPTX : 1 chunk par slide (texte + notes)

    # ── Plan du document (DOCX) / Sommaire slides (PPTX) ──────────────────
    outline: list[str] = field(default_factory=list)
    # Ex DOCX : ["1. Introduction", "2. État de l'art", ...]
    # Ex PPTX : ["Slide 1 : Contexte", "Slide 2 : Objectifs", ...]

    # ── Détail ─────────────────────────────────────────────────────────────
    blocks: list[DocxBlockResult] = field(default_factory=list)   # DOCX
    slides: list[SlideResult]     = field(default_factory=list)   # PPTX

    # ── Pages avec images/schémas à envoyer vers visual/ ──────────────────
    visual_candidates: list[int] = field(default_factory=list)
    # Numéros de slides/pages contenant des images

    # ── Métadonnées ────────────────────────────────────────────────────────
    metadata: DocumentMetadata = field(default_factory=DocumentMetadata)

    # ── Sections R&D détectées ─────────────────────────────────────────────
    detected_rd_sections: list[str] = field(default_factory=list)

    # ── Traçabilité ────────────────────────────────────────────────────────
    tags: list[str] = field(default_factory=list)
    confidence_score: float = 1.0
    extraction_errors: list[str] = field(default_factory=list)


# ── Utilitaires communs ───────────────────────────────────────────────────────


def _table_to_markdown(rows: list[list[str]]) -> str:
    """
    Convertit une liste de lignes/colonnes en tableau Markdown.
    Identique à pdf_native pour cohérence downstream (cleaner, RAG).
    """
    if not rows or not rows[0]:
        return ""

    cleaned = [
        [str(cell).strip().replace("\n", " ") if cell is not None else ""
         for cell in row]
        for row in rows
    ]

    if not cleaned:
        return ""

    col_count = max(len(r) for r in cleaned)
    col_widths = [
        max((len(r[i]) if i < len(r) else 0) for r in cleaned)
        for i in range(col_count)
    ]
    col_widths = [max(w, 3) for w in col_widths]

    def _fmt(row: list[str]) -> str:
        padded = [
            (row[i].ljust(col_widths[i]) if i < len(row) else " " * col_widths[i])
            for i in range(col_count)
        ]
        return "| " + " | ".join(padded) + " |"

    sep = "| " + " | ".join("-" * w for w in col_widths) + " |"
    lines = [_fmt(cleaned[0]), sep] + [_fmt(r) for r in cleaned[1:]]
    return "\n".join(lines)


def _detect_rd_sections(text: str) -> list[str]:
    matches = _RD_SECTION_RE.findall(text)
    sections = []
    for match in matches:
        if isinstance(match, tuple):
            sections.extend(s.strip() for s in match if s.strip())
        elif match.strip():
            sections.append(match.strip())
    seen: set[str] = set()
    unique: list[str] = []
    for s in sections:
        key = s.lower()
        if key not in seen:
            seen.add(key)
            unique.append(s)
    return unique


def _is_word_toc_or_field(raw_text: str, style_name: str = "") -> bool:
    """
    Ignore la table des matières Word.
    Sans ça, les sections CIR sont détectées dans le sommaire : PAGEREF / _Toc.
    """
    raw = str(raw_text or "").strip()
    low = raw.lower()
    style = str(style_name or "").lower().strip()

    if not raw:
        return True

    if style.startswith("toc"):
        return True

    if "pageref" in low or "_toc" in low or low.startswith("toc \\o"):
        return True

    if low in {"table des matières", "table des matieres", "sommaire"}:
        return True

    # Lignes de sommaire après extraction : 1.3. Etat de l'art 5
    if re.match(r"^\d+(?:\.\d+)*\.?\s+.{3,180}\s+\d{1,3}$", raw):
        words = re.findall(r"[A-Za-zÀ-ÿ]{3,}", raw)
        if len(words) >= 2:
            return True

    # Coordonnées Word parasites extraites depuis certains DOCX.
    if re.match(r"^-?\d{8,}$", raw):
        return True

    return False


# ══════════════════════════════════════════════════════════════════════════════
# DOCX
# ══════════════════════════════════════════════════════════════════════════════

def _docx_table_to_rows(table: "DocxTable") -> list[list[str]]:
    """Extrait les cellules d'un tableau Word."""
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def _extract_docx_comments(doc: "DocxDocument") -> list[str]:
    """
    Extrait les commentaires/annotations Word depuis le XML interne.
    Les commentaires sont souvent des notes de consultant très utiles pour le RAG.
    """
    comments: list[str] = []
    try:
        # Les commentaires sont dans word/comments.xml
        comments_part = doc.part.package.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
        )
        from lxml import etree
        ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        root = comments_part._element
        for comment_el in root.findall(f".//{{{ns}}}comment"):
            texts = [t.text for t in comment_el.iter(f"{{{ns}}}t") if t.text]
            full = " ".join(texts).strip()
            if full:
                comments.append(full)
    except Exception:
        pass    # Pas de commentaires ou XML inaccessible — non bloquant
    return comments


def _extract_docx_metadata(doc: "DocxDocument", first_text: str) -> DocumentMetadata:
    """Extrait les métadonnées d'un document Word."""
    props = doc.core_properties
    return DocumentMetadata(
        title=props.title or None,
        author=props.author or None,
        company=getattr(props, "company", None) or None,
        last_modified_by=props.last_modified_by or None,
        created=str(props.created) if props.created else None,
        modified=str(props.modified) if props.modified else None,
        revision=str(props.revision) if props.revision else None,
        organisme_detected=None,
        word_count=len(first_text.split()),
    )


def _build_docx_chunks(
    blocks: list[DocxBlockResult],
    comments: list[str],
) -> tuple[list[str], list[str]]:
    """
    Regroupe les blocs DOCX en chunks sémantiques pour le RAG.
    Intègre aussi les formules DANS chaque chunk (conserve le contexte).

    Logique de regroupement :
      - Un nouveau chunk commence à chaque Heading
      - Le contenu (paragraphes, tableaux) suit le heading courant
      - Les formules sont intégrées dans chaque chunk (pas comme chunks séparés)
      - Les commentaires forment un chunk dédié final

    Retourne (chunks, outline)
    """
    chunks: list[str] = []
    outline: list[str] = []
    current_heading: Optional[str] = None
    current_parts: list[str] = []

    def _flush() -> None:
        if current_parts:
            header = f"[SECTION : {current_heading}]\n" if current_heading else ""
            chunk = header + "\n\n".join(current_parts)
            # Formules intégrées par router.py
            chunks.append(chunk)

    for block in blocks:
        if block.block_type == "heading":
            _flush()
            current_heading = block.raw_text.strip()
            current_parts = []
            indent = "  " * max(0, block.level - 1)
            outline.append(f"{indent}{block.level}. {current_heading}")
        elif block.raw_text.strip():
            current_parts.append(block.raw_text.strip())

    _flush()    # Dernier groupe

    # Chunk commentaires consultant
    if comments:
        comment_block = "[COMMENTAIRES CONSULTANT]\n" + "\n---\n".join(comments)
        # Formules intégrées par router.py
        chunks.append(comment_block)

    return chunks, outline



def _xml_collect_text(element, ns: dict[str, str]) -> str:
    """
    Collecte le texte Word XML en respectant l'ordre des runs.
    Gère w:t, w:tab et w:br.
    """
    parts: list[str] = []

    for node in element.iter():
        tag = node.tag.split("}")[-1] if "}" in node.tag else node.tag

        if tag == "t" and node.text:
            parts.append(node.text)

        elif tag == "tab":
            parts.append("\t")

        elif tag in {"br", "cr"}:
            parts.append("\n")

    text = "".join(parts)
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_docx_xml_style_name(paragraph, ns: dict[str, str]) -> str:
    """
    Récupère le style XML d'un paragraphe Word.
    Exemple : Heading1, Titre1, Normal...
    """
    try:
        p_style = paragraph.find("./w:pPr/w:pStyle", ns)
        if p_style is None:
            return ""

        val = (
            p_style.attrib.get(f"{{{ns['w']}}}val")
            or p_style.attrib.get("w:val")
            or p_style.attrib.get("val")
            or ""
        )
        return str(val or "").lower().strip()

    except Exception:
        return ""


def _xml_table_to_rows(table_el, ns: dict[str, str]) -> list[list[str]]:
    """
    Extrait un tableau Word depuis le XML brut.
    """
    rows: list[list[str]] = []

    for tr in table_el.findall(".//w:tr", ns):
        row: list[str] = []

        for tc in tr.findall("./w:tc", ns):
            cell_text = _xml_collect_text(tc, ns)
            row.append(cell_text.replace("\n", " ").strip())

        if any(c.strip() for c in row):
            rows.append(row)

    return rows


def _extract_docx_xml_comments(zf: zipfile.ZipFile, ns: dict[str, str]) -> list[str]:
    """
    Extrait les commentaires Word depuis word/comments.xml si présent.
    """
    comments: list[str] = []

    try:
        if "word/comments.xml" not in zf.namelist():
            return comments

        root = ET.fromstring(zf.read("word/comments.xml"))

        for comment_el in root.findall(".//w:comment", ns):
            txt = _xml_collect_text(comment_el, ns)
            if txt:
                comments.append(txt)

    except Exception:
        pass

    return comments


def _extract_docm_xml_fallback(path: Path, file_type: str = OfficeFileType.DOCM.value) -> OfficeResult:
    """
    Fallback robuste pour DOCM/DOCX lorsque python-docx refuse le content type.
    Cas observé :
      application/vnd.ms-word.document.macroEnabled.main+xml

    Principe :
      - le .docm est un ZIP Office Open XML ;
      - on lit directement word/document.xml ;
      - on extrait paragraphes + tableaux dans l'ordre ;
      - on ignore les macros VBA.
    """
    result = OfficeResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_type=file_type,
    )

    ns = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    }

    try:
        if not zipfile.is_zipfile(str(path)):
            result.extraction_errors.append("Fallback XML impossible : fichier Office non ZIP")
            result.confidence_score = 0.10
            return result

        all_text_parts: list[str] = []
        block_index = 0

        with zipfile.ZipFile(str(path), "r") as zf:
            names = set(zf.namelist())

            xml_files: list[tuple[str, str]] = []

            if "word/document.xml" in names:
                xml_files.append(("document", "word/document.xml"))

            # Les en-têtes/pieds de page peuvent contenir des infos utiles.
            for n in sorted(names):
                if n.startswith("word/header") and n.endswith(".xml"):
                    xml_files.append(("header", n))

            for n in sorted(names):
                if n.startswith("word/footer") and n.endswith(".xml"):
                    xml_files.append(("footer", n))

            if not xml_files:
                result.extraction_errors.append(
                    "Fallback XML impossible : word/document.xml introuvable"
                )
                result.confidence_score = 0.10
                return result

            comments = _extract_docx_xml_comments(zf, ns)

            for xml_kind, xml_name in xml_files:
                try:
                    root = ET.fromstring(zf.read(xml_name))
                except Exception as exc:
                    result.extraction_errors.append(f"XML ignoré {xml_name}: {exc}")
                    continue

                body = root.find(".//w:body", ns)
                children = list(body) if body is not None else list(root)

                for child in children:
                    tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

                    if tag == "p":
                        raw_text = _xml_collect_text(child, ns)

                        if not raw_text or len(raw_text) < MIN_CHARS_PARA:
                            continue

                        style_name = _extract_docx_xml_style_name(child, ns)

                        if _is_word_toc_or_field(raw_text, style_name):
                            continue

                        style_low = style_name.lower()
                        is_heading = (
                            any(h.replace(" ", "") in style_low.replace(" ", "") for h in DOCX_HEADING_STYLES)
                            or style_low.startswith("heading")
                            or style_low.startswith("titre")
                        )

                        level = 0
                        if is_heading:
                            m = re.search(r"(\d+)", style_low)
                            level = int(m.group(1)) if m else 1

                        block = DocxBlockResult(
                            block_index=block_index,
                            block_type="heading" if is_heading else "paragraph",
                            style_name=style_name or f"xml_{xml_kind}",
                            raw_text=raw_text,
                            level=level,
                            is_rd_section=bool(_detect_rd_sections(raw_text)),
                            char_count=len(raw_text),
                        )

                        result.blocks.append(block)
                        all_text_parts.append(raw_text)
                        block_index += 1

                    elif tag == "tbl":
                        rows = _xml_table_to_rows(child, ns)
                        md = _table_to_markdown(rows)

                        if md:
                            block = DocxBlockResult(
                                block_index=block_index,
                                block_type="table",
                                style_name=f"xml_table_{xml_kind}",
                                raw_text=f"[TABLEAU]\n{md}",
                                level=0,
                                is_rd_section=False,
                                char_count=len(md),
                            )

                            result.blocks.append(block)
                            all_text_parts.append(f"[TABLEAU]\n{md}")
                            block_index += 1

        result.metadata = DocumentMetadata(
            title=None,
            author=None,
            word_count=sum(len(b.raw_text.split()) for b in result.blocks),
            has_comments=bool(comments),
        )

        result.text_chunks, result.outline = _build_docx_chunks(
            result.blocks,
            comments,
        )

        full_text = " ".join(all_text_parts)
        result.detected_rd_sections = _detect_rd_sections(full_text)

        result.tags = [
            file_type.upper(),
            "XML_FALLBACK",
        ]

        if file_type == OfficeFileType.DOCM.value:
            result.tags.append("MACRO_ENABLED_DOCM")

        if any(b.block_type == "table" for b in result.blocks):
            result.tags.append("HAS_TABLES")

        if comments:
            result.tags.append("HAS_COMMENTS")

        if result.outline:
            result.tags.append("HAS_STRUCTURE")

        if result.detected_rd_sections:
            result.tags.append("CIR_SECTIONS")

        if not result.text_chunks:
            result.extraction_errors.append(
                "Fallback XML : aucun texte exploitable extrait"
            )
            result.confidence_score = 0.20
        else:
            result.confidence_score = 0.75

        return result

    except Exception as exc:
        result.extraction_errors.append(f"Fallback XML impossible : {exc}")
        result.confidence_score = 0.10
        return result


def _extract_docx(path: Path, file_type: str = OfficeFileType.DOCX.value) -> OfficeResult:
    """Pipeline complet d'extraction DOCX/DOCM."""
    result = OfficeResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_type=file_type,
    )

    if not DOCX_AVAILABLE:
        result.extraction_errors.append(
            "python-docx non installé : pip install python-docx"
        )
        return result

    try:
        doc = DocxDocument(str(path))
    except Exception as exc:
        msg = str(exc)

        # Le XML principal peut rester parfaitement exploitable même si
        # python-docx refuse le paquet (DOCM, média annexe au CRC invalide,
        # relation non standard, etc.). Le fallback ne lit que les parties XML
        # textuelles et évite qu'une image corrompue annule tout le CIR.
        fallback = _extract_docm_xml_fallback(path, file_type=file_type)

        if fallback.text_chunks:
            fallback.tags.append("PYTHON_DOCX_FALLBACK_USED")
            if path.suffix.lower() == ".docm" or "macroEnabled" in msg or "macroEnabled.main+xml" in msg:
                fallback.tags.append("DOCM_CONTENT_TYPE_FALLBACK")
            else:
                fallback.tags.append("DOCX_PARTIAL_ZIP_FALLBACK")
            return fallback

        fallback.extraction_errors.append(
            f"python-docx a refusé le document et le fallback XML n'a pas extrait de chunks : {exc}"
        )
        return fallback

    all_text_parts: list[str] = []
    block_index = 0

    # ── Parcours du corps du document ─────────────────────────────────────
    # doc.element.body contient paragraphes ET tableaux dans l'ordre réel
    from docx.oxml.ns import qn
    for child in doc.element.body:

        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        # ── Paragraphe ────────────────────────────────────────────────────
        if tag == "p":
            para = DocxParagraph(child, doc)
            raw_text = para.text.strip()
            if not raw_text or len(raw_text) < MIN_CHARS_PARA:
                continue
            style_obj = getattr(para, "style", None)
            style_name = (getattr(style_obj, "name", None) or "").lower()

            # Important : ne pas envoyer la table des matières Word au NLP/CIR.
            if _is_word_toc_or_field(raw_text, style_name):
                continue

            is_heading = any(h in style_name for h in DOCX_HEADING_STYLES)
            level = 0
            if is_heading:
                # Extraire le niveau du style (ex: "Heading 2" → 2)
                m = re.search(r"(\d+)", style_name)
                level = int(m.group(1)) if m else 1

            block = DocxBlockResult(
                block_index=block_index,
                block_type="heading" if is_heading else "paragraph",
                style_name=style_name,
                raw_text=raw_text,
                level=level,
                is_rd_section=bool(_detect_rd_sections(raw_text)),
                char_count=len(raw_text),
            )
            result.blocks.append(block)
            all_text_parts.append(raw_text)
            block_index += 1

        # ── Tableau ───────────────────────────────────────────────────────
        elif tag == "tbl":
            from docx.table import Table as _Table
            table = _Table(child, doc)
            rows = _docx_table_to_rows(table)
            md = _table_to_markdown(rows)
            if md:
                block = DocxBlockResult(
                    block_index=block_index,
                    block_type="table",
                    style_name="table",
                    raw_text=f"[TABLEAU]\n{md}",
                    level=0,
                    is_rd_section=False,
                    char_count=len(md),
                )
                result.blocks.append(block)
                all_text_parts.append(f"[TABLEAU]\n{md}")
                block_index += 1

    # ── Commentaires ──────────────────────────────────────────────────────
    comments = _extract_docx_comments(doc)
    result.metadata.has_comments = bool(comments)

    # ── Métadonnées ───────────────────────────────────────────────────────
    first_text = " ".join(all_text_parts[:5])
    result.metadata = _extract_docx_metadata(doc, " ".join(all_text_parts))
    result.metadata.has_comments = bool(comments)
    result.metadata.word_count = sum(
        len(b.raw_text.split()) for b in result.blocks
    )

    # ── Construction chunks + outline ──────────────────────────────────────────────
    # Les formules sont maintenant intégrées DANS les chunks via _build_docx_chunks()
    result.text_chunks, result.outline = _build_docx_chunks(result.blocks, comments)

    # ── Sections R&D globales ─────────────────────────────────────────────
    full_text = " ".join(all_text_parts)
    result.detected_rd_sections = _detect_rd_sections(full_text)

    return result


# ══════════════════════════════════════════════════════════════════════════════
# PPTX
# ══════════════════════════════════════════════════════════════════════════════

def _extract_pptx_metadata(prs: "Presentation", first_text: str) -> DocumentMetadata:
    """Extrait les métadonnées d'une présentation PowerPoint."""
    props = prs.core_properties
    return DocumentMetadata(
        title=props.title or None,
        author=props.author or None,
        company=getattr(props, "company", None) or None,
        last_modified_by=props.last_modified_by or None,
        created=str(props.created) if props.created else None,
        modified=str(props.modified) if props.modified else None,
        organisme_detected=None,
        slide_count=0,   # Mis à jour après
    )


def _extract_slide_title(slide: "Presentation") -> Optional[str]:
    """Extrait le titre d'une slide (placeholder title ou premier texte)."""
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    # Fallback : premier texte non vide
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text[:80]
    return None


def _extract_slide(slide: object, slide_number: int) -> SlideResult:
    """
    Extrait tout le contenu d'une slide PPTX avec support récursif profond.
    Gère les Groupes, les Tableaux, les GraphicFrames (Excel/SmartArt) et les Formes.
    """
    title = _extract_slide_title(slide)
    text_parts: list[str] = []
    tables_as_text: list[str] = []
    has_images = False
    has_charts = False

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _process_shapes(shapes):
        nonlocal has_images, has_charts
        for shape in shapes:
            # On récupère le type de forme de manière sécurisée (ID numérique)
            # 6  = GROUP
            # 12 = GRAPHIC_FRAME (Tableaux Excel, SmartArt, Graphiques)
            # 13 = PICTURE
            stype = getattr(shape, "shape_type", -1)

            # 1. GESTION DES GROUPES (Recursivité) - ID 6
            if stype == 6:
                try:
                    _process_shapes(shape.shapes)
                except:
                    pass
                continue

            # 2. GESTION DES TABLEAUX DIRECTS (Le cas standard)
            if shape.has_table:
                try:
                    rows = []
                    for row in shape.table.rows:
                        rows.append([
                            cell.text_frame.text.strip().replace("\n", " ") 
                            if cell.text_frame else ""
                            for cell in row.cells
                        ])
                    md = _table_to_markdown(rows)
                    if md: tables_as_text.append(md)
                    continue # On passe à la forme suivante pour éviter les doublons
                except:
                    pass

            # 3. GESTION DES GRAPHIC FRAMES (Tableaux Excel / SmartArt) - ID 12
            # C'est ici que se cache le tableau de ta Slide 11
            if stype == 12:
                # A. On cherche si un tableau est imbriqué dans le frame
                if hasattr(shape, "has_table") and shape.has_table:
                    try:
                        rows = []
                        for row in shape.table.rows:
                            rows.append([
                                c.text_frame.text.strip().replace("\n", " ") 
                                for c in row.cells
                            ])
                        md = _table_to_markdown(rows)
                        if md: tables_as_text.append(md)
                    except:
                        pass
                
                # B. On cherche si c'est un graphique
                if hasattr(shape, "has_chart") and shape.has_chart:
                    has_charts = True
                    try:
                        chart_title = shape.chart.chart_title.text_frame.text
                        text_parts.append(f"[GRAPHIQUE : {chart_title}]")
                    except:
                        text_parts.append("[GRAPHIQUE]")
                continue

            # 4. GESTION DU TEXTE (Zones de texte et Formes simples)
            if shape.has_text_frame:
                try:
                    shape_text = shape.text_frame.text.strip()
                    if shape_text and shape_text != (title or ""):
                        # Détection des tableaux "dessinés" à la main (tabulations)
                        if "\t" in shape_text or "    " in shape_text:
                            text_parts.append(f"[DONNÉES TABULAIRES DÉTECTÉES]\n{shape_text}")
                        else:
                            text_parts.append(shape_text)
                except:
                    pass

            # 5. IMAGES (ID 13)
            if stype == 13:
                has_images = True
    # Lancement de l'extraction sur la collection de formes de la slide
    _process_shapes(slide.shapes)

    # 6. NOTES PRÉSENTATEUR (Source riche pour le RAG)
    notes_text = ""
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
    except Exception:
        pass

    # Assemblage final
    raw_text = "\n".join(text_parts)
    full_content = raw_text + " " + notes_text
    char_count = len(raw_text.strip())

    return SlideResult(
        slide_number=slide_number,
        title=title,
        raw_text=raw_text,
        notes_text=notes_text,
        tables_as_text=tables_as_text,
        has_images=has_images,
        has_charts=has_charts,
        detected_sections=_detect_rd_sections(full_content),
        char_count=char_count,
        is_text_poor=char_count < MIN_CHARS_SLIDE,
    )

def _build_slide_chunk(slide: SlideResult) -> str:
    """
    Assemble le chunk RAG d'une slide avec les formules intégrées.

    Format :
        [SLIDE 3 : Titre de la slide]
        <texte des zones>

        [NOTES PRÉSENTATEUR]
        <notes>

        [TABLEAU 1]
        | col1 | col2 |
        
        [FORMULES DÉTECTÉES]
        1. LaTeX: ... 
           Explication: ...
           Domaine: ...
           Confiance: ...%
    """
    title_str = f" : {slide.title}" if slide.title else ""
    parts = [f"[SLIDE {slide.slide_number}{title_str}]"]

    if slide.raw_text.strip():
        parts.append(slide.raw_text.strip())

    if slide.notes_text.strip():
        parts.append(f"[NOTES PRÉSENTATEUR]\n{slide.notes_text.strip()}")

    for i, table_md in enumerate(slide.tables_as_text, start=1):
        parts.append(f"[TABLEAU {i}]\n{table_md}")

    # Assemblage avant intégration des formules et images
    chunk = "\n\n".join(parts)
    
    # Formules et images intégrées par router.py
    
    return chunk


def _extract_pptx(path: Path) -> OfficeResult:
    """Pipeline complet d'extraction PPTX."""
    result = OfficeResult(
        file_name=path.name,
        source_path=str(path.resolve()),
        file_type=OfficeFileType.PPTX.value,
    )

    if not PPTX_AVAILABLE:
        result.extraction_errors.append(
            "python-pptx non installé : pip install python-pptx"
        )
        return result

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        result.extraction_errors.append(f"Impossible d'ouvrir le PPTX : {exc}")
        return result

    all_text_parts: list[str] = []
    all_rd_sections: set[str] = set()

    for i, slide in enumerate(prs.slides, start=1):
        try:
            slide_result = _extract_slide(slide, i)
            result.slides.append(slide_result)
            # Le chunk maintenant inclut les formules intégrées (via _build_slide_chunk)
            result.text_chunks.append(_build_slide_chunk(slide_result))

            # Outline
            slide_title = slide_result.title or f"Slide {i}"
            result.outline.append(f"Slide {i} : {slide_title}")

            # Candidats visuels
            if slide_result.has_images or slide_result.has_charts:
                result.visual_candidates.append(i)

            # Sections R&D
            all_rd_sections.update(slide_result.detected_sections)

            # Texte global
            all_text_parts.append(slide_result.raw_text)
            all_text_parts.append(slide_result.notes_text)

        except Exception as exc:
            msg = f"Erreur slide {i} : {exc}"
            logger.warning(msg)
            result.extraction_errors.append(msg)
            result.text_chunks.append(f"[SLIDE {i}]\n[ERREUR EXTRACTION]")

    # ── Métadonnées ───────────────────────────────────────────────────────
    full_text = " ".join(all_text_parts)
    result.metadata = _extract_pptx_metadata(prs, full_text[:2000])
    result.metadata.slide_count = len(prs.slides)

    result.detected_rd_sections = sorted(all_rd_sections)

    return result


# ── Tags & Confiance ──────────────────────────────────────────────────────────

def _build_tags(result: OfficeResult) -> list[str]:
    """
    Construit les tags sans écraser les tags déjà posés par les fallback.
    Exemple important : XML_FALLBACK / MACRO_ENABLED_DOCM / PYTHON_DOCX_FALLBACK_USED.
    """
    tags: list[str] = list(result.tags or [])

    main_type = str(result.file_type or "").upper()
    if main_type and main_type not in tags:
        tags.insert(0, main_type)

    if result.file_type in {OfficeFileType.PPTX.value, OfficeFileType.PPTM.value}:
        if result.visual_candidates:
            tags.append("HAS_VISUAL_SLIDES")
        if any(s.notes_text for s in result.slides):
            tags.append("HAS_PRESENTER_NOTES")
        if any(s.tables_as_text for s in result.slides):
            tags.append("HAS_TABLES")

    elif result.file_type in {OfficeFileType.DOCX.value, OfficeFileType.DOCM.value}:
        if any(b.block_type == "table" for b in result.blocks):
            tags.append("HAS_TABLES")
        if result.metadata.has_comments:
            tags.append("HAS_COMMENTS")
        if result.outline:
            tags.append("HAS_STRUCTURE")

    if result.detected_rd_sections:
        tags.append("CIR_SECTIONS")

    # PARTIAL_EXTRACTION doit seulement représenter une vraie erreur restante.
    if result.extraction_errors:
        tags.append("PARTIAL_EXTRACTION")

    return list(dict.fromkeys(tags))


def _compute_confidence(result: OfficeResult) -> float:
    """
    Score de confiance basé sur la richesse du contenu extrait.

    DOCX : pénalise les blocs vides et les erreurs
    PPTX : pénalise les slides pauvres et les erreurs
    """
    if not result.text_chunks:
        return 0.0

    error_penalty = min(len(result.extraction_errors) * 0.10, 0.30)

    if result.file_type in {OfficeFileType.PPTX.value, OfficeFileType.PPTM.value} and result.slides:
        poor_ratio = sum(1 for s in result.slides if s.is_text_poor) / len(result.slides)
        score = 1.0 - (poor_ratio * 0.40) - error_penalty
    else:
        score = 1.0 - error_penalty

    return max(round(score, 2), 0.10)


# ── Point d'entrée principal ──────────────────────────────────────────────────

def extract_office(file_path: str | Path) -> OfficeResult:
    """
    Extrait le contenu d'un document Office (DOCX/DOCM ou PPTX/PPTM) pour le RAG EnnoSmart.

    Paramètres
    ----------
    file_path : str | Path
        Chemin vers le fichier .docx, .docm, .pptx ou .pptm

    Retourne
    --------
    OfficeResult
        text_chunks          : chunks RAG (1 par section DOCX / 1 par slide PPTX)
        outline              : plan du document / sommaire des slides
        visual_candidates    : numéros de slides avec images → à envoyer vers vision.py
        metadata             : auteur, société, organisme détecté
        detected_rd_sections : sections CIR trouvées
        tags                 : traçabilité
        confidence_score     : qualité globale

    Raises
    ------
    FileNotFoundError : fichier introuvable
    ValueError        : extension non supportée (.doc, .ppt non supportés sans conversion)
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"Fichier introuvable : {path}")

    ext = path.suffix.lower()

    # .docm et .pptm sont des formats Office Open XML macro-enabled.
    # Ils ont une structure ZIP proche de .docx/.pptx et peuvent être lus
    # par python-docx / python-pptx dans la majorité des cas.
    supported_word = {".docx", ".docm"}
    supported_powerpoint = {".pptx", ".pptm"}
    supported = supported_word | supported_powerpoint

    if ext not in supported:
        raise ValueError(
            f"Extension '{ext}' non supportée par office.py. "
            f"Formats acceptés : .docx, .docm, .pptx, .pptm\n"
            f"Note : les anciens formats .doc/.ppt nécessitent une conversion préalable."
        )

    logger.info("Extraction Office [%s] → %s", ext.upper(), path.name)

    if ext in supported_word:
        result = _extract_docx(path, file_type=ext.lstrip("."))
    else:
        result = _extract_pptx(path)
        result.file_type = ext.lstrip(".")

    result.tags = _build_tags(result)
    result.confidence_score = _compute_confidence(result)

    logger.info(
        "✓ %s — %d chunks | score=%.2f | tags=%s | rd_sections=%d",
        path.name,
        len(result.text_chunks),
        result.confidence_score,
        result.tags,
        len(result.detected_rd_sections),
    )

    return result


# ── Interface rapide (debug / tests) ──────────────────────────────────────────

if __name__ == "__main__":
    import sys
    import json

    logging.basicConfig(level=logging.DEBUG)

    if len(sys.argv) < 2:
        print("Usage : python office.py <chemin_vers_fichier.docx|.docm|.pptx|.pptm>")
        sys.exit(1)

    res = extract_office(sys.argv[1])

    summary = {
        "file":           res.file_name,
        "type":           res.file_type,
        "chunks":         len(res.text_chunks),
        "outline":        res.outline[:10],
        "organisme":      res.metadata.organisme_detected,
        "rd_sections":    res.detected_rd_sections,
        "visual_slides":  res.visual_candidates,
        "tags":           res.tags,
        "confidence":     res.confidence_score,
        "errors":         res.extraction_errors,
        "chunks_preview": [c[:300] + "…" for c in res.text_chunks[:3]],
    }

    print(json.dumps(summary, ensure_ascii=False, indent=2))
