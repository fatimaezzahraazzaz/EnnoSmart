# -*- coding: utf-8 -*-
from __future__ import annotations

"""Conservation structurelle des images, sans modèle de vision.

Les fichiers extraits sont rattachés à leur paragraphe, slide, feuille ou page
dans un manifeste séparé. Ce manifeste n'est pas injecté dans le NLP/RAG : le
diagnostic continue à s'appuyer sur le texte et les tableaux structurés.
"""

import hashlib
import json
import posixpath
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Iterable, Mapping, Optional


VERSION = "visual_asset_retention_v1_no_vision_model"
DIRNAME = "visual_assets_v1"
STANDALONE_IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".tif",
    ".tiff",
    ".bmp",
    ".gif",
    ".webp",
    ".svg",
}


def _safe_name(value: Any, default: str = "document") -> str:
    name = Path(str(value or default)).name
    cleaned = re.sub(r"[^A-Za-z0-9À-ÿ._-]+", "_", name).strip("._")
    return cleaned[:140] or default


def _context(value: Any, limit: int = 1200) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text[:limit]


def _write_asset(
    target_dir: Path,
    data: bytes,
    extension: str,
    cache: dict[str, str],
) -> str:
    digest = hashlib.sha256(data).hexdigest()
    if digest in cache:
        return cache[digest]

    ext = re.sub(r"[^a-z0-9]", "", str(extension or "").lower()) or "bin"
    path = target_dir / f"{digest[:20]}.{ext}"
    if not path.exists():
        path.write_bytes(data)
    cache[digest] = str(path.resolve())
    return cache[digest]


def _docx_assets(
    source: Path,
    target: Path,
    document: str,
    cache: dict[str, str],
) -> list[dict[str, Any]]:
    namespaces = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
    }
    records: list[dict[str, Any]] = []
    with zipfile.ZipFile(str(source)) as archive:
        document_xml = ET.fromstring(archive.read("word/document.xml"))
        rels_xml = ET.fromstring(archive.read("word/_rels/document.xml.rels"))
        relationships = {
            str(rel.attrib.get("Id") or ""): str(rel.attrib.get("Target") or "")
            for rel in rels_xml.findall("pr:Relationship", namespaces)
        }
        paragraphs = document_xml.findall(".//w:p", namespaces)
        paragraph_texts = [
            _context(" ".join(node.text or "" for node in p.findall(".//w:t", namespaces)))
            for p in paragraphs
        ]

        for paragraph_index, paragraph in enumerate(paragraphs):
            embeds = [
                str(node.attrib.get(f"{{{namespaces['r']}}}embed") or "")
                for node in paragraph.findall(".//a:blip", namespaces)
            ]
            embeds = [value for value in embeds if value]
            if not embeds:
                continue
            nearby = [
                paragraph_texts[index]
                for index in (paragraph_index - 1, paragraph_index, paragraph_index + 1)
                if 0 <= index < len(paragraph_texts) and paragraph_texts[index]
            ]
            for relation_id in embeds:
                relation_target = relationships.get(relation_id, "")
                member = posixpath.normpath(posixpath.join("word", relation_target))
                if member not in archive.namelist():
                    continue
                data = archive.read(member)
                asset_path = _write_asset(
                    target,
                    data,
                    Path(member).suffix,
                    cache,
                )
                records.append(
                    {
                        "document": document,
                        "asset_path": asset_path,
                        "anchor_type": "docx_paragraph",
                        "paragraph_index": paragraph_index,
                        "paragraph_context": _context(" ".join(nearby)),
                        "relationship_id": relation_id,
                        "original_member": member,
                    }
                )
    return records


def _pptx_assets(
    source: Path,
    target: Path,
    document: str,
    cache: dict[str, str],
) -> list[dict[str, Any]]:
    from pptx import Presentation

    records: list[dict[str, Any]] = []
    presentation = Presentation(str(source))

    def walk(shapes: Any) -> list[Any]:
        output: list[Any] = []
        for shape in shapes:
            output.append(shape)
            nested = getattr(shape, "shapes", None)
            if nested is not None:
                output.extend(walk(nested))
        return output

    for slide_number, slide in enumerate(presentation.slides, start=1):
        shapes = walk(slide.shapes)
        texts = []
        for shape in shapes:
            text = _context(getattr(shape, "text", ""))
            if not text:
                continue
            center = (
                float(getattr(shape, "left", 0)) + float(getattr(shape, "width", 0)) / 2,
                float(getattr(shape, "top", 0)) + float(getattr(shape, "height", 0)) / 2,
            )
            texts.append((center, text))

        for shape_index, shape in enumerate(shapes, start=1):
            image = getattr(shape, "image", None)
            if image is None:
                continue
            data = bytes(image.blob)
            asset_path = _write_asset(
                target,
                data,
                getattr(image, "ext", "bin"),
                cache,
            )
            center = (
                float(getattr(shape, "left", 0)) + float(getattr(shape, "width", 0)) / 2,
                float(getattr(shape, "top", 0)) + float(getattr(shape, "height", 0)) / 2,
            )
            nearest_text = ""
            if texts:
                nearest_text = min(
                    texts,
                    key=lambda item: (
                        (item[0][0] - center[0]) ** 2
                        + (item[0][1] - center[1]) ** 2
                    ),
                )[1]
            records.append(
                {
                    "document": document,
                    "asset_path": asset_path,
                    "anchor_type": "pptx_slide",
                    "slide_number": slide_number,
                    "shape_index": shape_index,
                    "paragraph_context": nearest_text,
                }
            )
    return records


def _excel_assets(
    source: Path,
    target: Path,
    document: str,
    cache: dict[str, str],
) -> list[dict[str, Any]]:
    import openpyxl

    records: list[dict[str, Any]] = []
    workbook = openpyxl.load_workbook(str(source), data_only=True, read_only=False)
    try:
        for worksheet in workbook.worksheets:
            for image_index, image in enumerate(
                getattr(worksheet, "_images", None) or [],
                start=1,
            ):
                anchor = getattr(image, "anchor", None)
                marker = getattr(anchor, "_from", None)
                row = int(getattr(marker, "row", 0)) + 1
                column = int(getattr(marker, "col", 0)) + 1
                nearby_values: list[str] = []
                for row_index in range(max(1, row - 1), row + 2):
                    values = [
                        str(cell.value)
                        for cell in worksheet[row_index]
                        if cell.value not in (None, "")
                    ]
                    if values:
                        nearby_values.append(" | ".join(values))
                data = bytes(image._data())
                asset_path = _write_asset(
                    target,
                    data,
                    getattr(image, "format", "bin"),
                    cache,
                )
                records.append(
                    {
                        "document": document,
                        "asset_path": asset_path,
                        "anchor_type": "excel_cell",
                        "sheet_name": worksheet.title,
                        "row": row,
                        "column": column,
                        "image_index": image_index,
                        "paragraph_context": _context(" ".join(nearby_values)),
                    }
                )
    finally:
        workbook.close()
    return records


def _pdf_assets(
    source: Path,
    target: Path,
    document: str,
    cache: dict[str, str],
) -> list[dict[str, Any]]:
    import fitz

    records: list[dict[str, Any]] = []
    pdf = fitz.open(str(source))
    try:
        for page_index in range(len(pdf)):
            page = pdf[page_index]
            blocks = page.get_text("blocks") or []
            for image_index, image_info in enumerate(page.get_images(full=True), start=1):
                xref = int(image_info[0])
                extracted = pdf.extract_image(xref)
                data = bytes(extracted.get("image") or b"")
                if not data:
                    continue
                rects = page.get_image_rects(xref)
                nearest_text = ""
                if rects and blocks:
                    rect = rects[0]
                    center_x = (rect.x0 + rect.x1) / 2
                    center_y = (rect.y0 + rect.y1) / 2
                    nearest = min(
                        blocks,
                        key=lambda block: (
                            (((float(block[0]) + float(block[2])) / 2) - center_x) ** 2
                            + (((float(block[1]) + float(block[3])) / 2) - center_y) ** 2
                        ),
                    )
                    nearest_text = _context(nearest[4])
                asset_path = _write_asset(
                    target,
                    data,
                    extracted.get("ext") or "bin",
                    cache,
                )
                records.append(
                    {
                        "document": document,
                        "asset_path": asset_path,
                        "anchor_type": "pdf_page",
                        "page_number": page_index + 1,
                        "image_index": image_index,
                        "paragraph_context": nearest_text,
                    }
                )
    finally:
        pdf.close()
    return records


def _standalone_image_assets(
    source: Path,
    target: Path,
    document: str,
    cache: dict[str, str],
) -> list[dict[str, Any]]:
    """Conserve une image déposée seule, sans OCR ni description générée."""
    asset_path = _write_asset(
        target,
        source.read_bytes(),
        source.suffix,
        cache,
    )
    return [
        {
            "document": document,
            "asset_path": asset_path,
            "anchor_type": "standalone_image",
            "paragraph_context": "",
        }
    ]


def persist_visual_assets(
    *,
    store: Any,
    documents: Iterable[Mapping[str, Any]],
) -> dict[str, Any]:
    root = Path(store.documents_processed_dir) / DIRNAME
    root.mkdir(parents=True, exist_ok=True)

    records: list[dict[str, Any]] = []
    errors: list[dict[str, str]] = []
    cache: dict[str, str] = {}
    handlers = {
        ".docx": _docx_assets,
        ".docm": _docx_assets,
        ".pptx": _pptx_assets,
        ".pptm": _pptx_assets,
        ".xlsx": _excel_assets,
        ".xlsm": _excel_assets,
        ".pdf": _pdf_assets,
    }

    for index, raw in enumerate(documents or [], start=1):
        source = Path(str(raw.get("source_path") or ""))
        if not source.exists() or not source.is_file():
            continue
        suffix = source.suffix.lower()
        handler = handlers.get(suffix)
        if handler is None and suffix in STANDALONE_IMAGE_EXTENSIONS:
            handler = _standalone_image_assets
        if handler is None:
            continue
        document = _safe_name(raw.get("document") or source.name, f"document_{index}")
        document_dir = root / f"{index:03d}_{Path(document).stem}"
        document_dir.mkdir(parents=True, exist_ok=True)
        try:
            extracted = handler(source, document_dir, document, cache)
            for record in extracted:
                record.update(
                    {
                        "source_path": str(source.resolve()),
                        "indexed_for_diagnostic": False,
                        "vision_model_used": False,
                    }
                )
            records.extend(extracted)
        except Exception as exc:
            errors.append({"document": document, "error": str(exc)})

    manifest = {
        "version": VERSION,
        "vision_model_used": False,
        "indexed_for_diagnostic": False,
        "assets_count": len(records),
        "unique_files_count": len(cache),
        "errors": errors,
        "assets": records,
    }
    manifest_path = root / "manifest.json"
    manifest_path.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return {
        **manifest,
        "manifest_path": str(manifest_path.resolve()),
    }
